from fastapi import FastAPI, APIRouter, HTTPException, Depends, status, UploadFile, File, Query
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from fastapi.responses import FileResponse, StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import re
from pathlib import Path
from pydantic import BaseModel, Field, EmailStr, HttpUrl
from typing import List, Optional, Literal
import uuid
from datetime import datetime, timezone, timedelta
import jwt
import bcrypt
from openai import OpenAI
import PyPDF2
import io
import aiofiles
import httpx
from bs4 import BeautifulSoup
from docx import Document
import pytesseract
from PIL import Image

# Enterprise Knowledge Architecture imports
from services.enterprise import AuditService, VersionService, HierarchicalRetrieval
from routes.departments import setup_department_routes, router as departments_router
from routes.enterprise_sources import setup_enterprise_source_routes, router as enterprise_sources_router
from routes.news import router as news_router
from routes.analyzer import setup_analyzer_routes, router as analyzer_router

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# JWT Settings
JWT_SECRET = os.environ.get('JWT_SECRET', 'shared-project-gpt-secret-key-2024')
JWT_ALGORITHM = "HS256"
JWT_EXPIRATION_HOURS = 24

# Admin email domain
ADMIN_EMAIL_DOMAIN = "@admin.com"

# OpenAI API Key
OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY', '')

# Claude API Key
CLAUDE_API_KEY = os.environ.get('CLAUDE_API_KEY', '')

# Initialize OpenAI client (for embeddings only)
openai_client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None

# File storage settings
UPLOAD_DIR = ROOT_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)
IMAGES_DIR = ROOT_DIR / "generated_images"
IMAGES_DIR.mkdir(exist_ok=True)
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
CHUNK_SIZE = 1500  # characters per chunk
CHUNK_SIZE_TABULAR = 800  # smaller chunks for tabular data (Excel, CSV)
MAX_CONTEXT_CHARS = 10000  # Max characters to include in context (reduced for token optimization)
MAX_CHUNKS_PER_QUERY = 5  # Max chunks to include per query (top 5 most relevant)
MAX_AUTO_INGEST_URLS = 3  # Max URLs to auto-ingest per message

# Global sources project ID marker
GLOBAL_PROJECT_ID = "__global__"  # Special marker for global sources

# Semantic cache settings
CACHE_SIMILARITY_THRESHOLD = 0.92  # Minimum cosine similarity to consider a cache hit
CACHE_TTL_DAYS = 30  # Cache entries expire after this many days
EMBEDDING_MODEL = "text-embedding-3-small"  # OpenAI embedding model

# Image generation settings
IMAGE_RATE_LIMIT_PER_HOUR = 10  # Max images per user per hour
VALID_IMAGE_SIZES = ["1024x1024", "1024x1792", "1792x1024"]

# URL pattern for detecting URLs in messages
URL_PATTERN = re.compile(r'https?://[^\s<>"{}|\\^`\[\]]+', re.IGNORECASE)

# Supported MIME types
SUPPORTED_MIME_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "text/plain": "txt",
    "text/markdown": "md",
    "text/csv": "csv",
    "application/csv": "csv",
    "image/png": "png",
    "image/jpeg": "jpeg",
    "image/jpg": "jpg",
}

# Create the main app
app = FastAPI(title="Shared Project GPT")

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Security
security = HTTPBearer()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ==================== PAGINATION HELPER ====================

class PaginatedResponse(BaseModel):
    items: List[dict]
    total: int
    page: int
    pageSize: int
    totalPages: int

async def paginate_query(collection, query: dict, page: int = 1, page_size: int = 50, 
                         sort_field: str = "createdAt", sort_order: int = -1, 
                         projection: dict = None):
    """Generic pagination helper for MongoDB queries"""
    if projection is None:
        projection = {"_id": 0}
    skip = (page - 1) * page_size
    total = await collection.count_documents(query)
    items = await collection.find(query, projection).sort(sort_field, sort_order).skip(skip).limit(page_size).to_list(page_size)
    total_pages = (total + page_size - 1) // page_size if total > 0 else 1
    return {
        "items": items,
        "total": total,
        "page": page,
        "pageSize": page_size,
        "totalPages": total_pages
    }

# ==================== MODELS ====================

class UserCreate(BaseModel):
    email: EmailStr
    password: str

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class UserResponse(BaseModel):
    id: str
    email: str
    isAdmin: bool
    createdAt: str
    canEditGlobalSources: Optional[bool] = False
    # Enterprise Knowledge Architecture
    departments: Optional[List[str]] = []
    primaryDepartmentId: Optional[str] = None

class UserWithUsageResponse(BaseModel):
    id: str
    email: str
    isAdmin: bool
    createdAt: str
    totalTokensUsed: int
    totalMessagesCount: int
    canEditGlobalSources: Optional[bool] = False

class UpdateUserGlobalPermissionRequest(BaseModel):
    canEditGlobalSources: bool

class TokenResponse(BaseModel):
    token: str
    user: UserResponse

class ProjectCreate(BaseModel):
    name: str

class ProjectRole(str):
    """Project sharing roles"""
    VIEWER = "viewer"      # Read-only: view chats and sources
    EDITOR = "editor"      # Create/edit chats, but no source management
    MANAGER = "manager"    # Full access: manage sources and members

class ProjectMember(BaseModel):
    userId: str
    email: str
    role: str  # viewer, editor, manager

class ProjectResponse(BaseModel):
    id: str
    name: str
    ownerId: str
    sharedWith: Optional[List[str]] = []
    sharedMembers: Optional[List[ProjectMember]] = []  # With roles
    createdAt: str

class ShareProjectRequest(BaseModel):
    email: str
    role: Optional[str] = "viewer"  # Default to viewer

class ChatCreate(BaseModel):
    name: Optional[str] = "New Chat"

class ChatResponse(BaseModel):
    id: str
    projectId: Optional[str] = None  # None for quick chats
    name: str
    createdAt: str
    activeSourceIds: Optional[List[str]] = []
    sharedWithUsers: Optional[List[str]] = None  # None = visible to all shared users, [] or [ids] = only those users
    sourceMode: Optional[str] = "all"  # 'all' or 'my'

class QuickChatCreate(BaseModel):
    name: Optional[str] = "Quick Chat"

class MoveChatRequest(BaseModel):
    targetProjectId: str

class RenameChatRequest(BaseModel):
    name: str

class UpdateChatVisibilityRequest(BaseModel):
    sharedWithUsers: List[str]  # List of user IDs who can see this chat

class MessageCreate(BaseModel):
    content: str

class EnhancedCitation(BaseModel):
    """Enhanced citation with full context"""
    sourceId: str
    sourceName: str
    sourceType: str  # "project" or "global"
    chunkId: str
    chunkIndex: int
    textFragment: str  # First 200 chars of chunk
    score: float

class MessageResponse(BaseModel):
    id: str
    chatId: str
    role: Literal["user", "assistant"]
    content: str
    createdAt: str
    citations: Optional[List[dict]] = None
    usedSources: Optional[List[dict]] = None  # For UI to reliably show "Sources used"
    autoIngestedUrls: Optional[List[str]] = None  # IDs of auto-ingested URL sources
    senderEmail: Optional[str] = None  # Email of user who sent the message
    senderName: Optional[str] = None  # Display name of sender
    fromCache: Optional[bool] = False  # Whether response is from cache
    cacheInfo: Optional[dict] = None  # Cache hit details

class GPTConfigUpdate(BaseModel):
    model: Optional[str] = None
    developerPrompt: Optional[str] = None

class GPTConfigResponse(BaseModel):
    id: str
    model: str
    developerPrompt: str
    updatedAt: str

class SourceResponse(BaseModel):
    id: str
    projectId: str
    kind: Literal["file", "url", "knowledge"]
    originalName: Optional[str] = None
    url: Optional[str] = None
    mimeType: Optional[str] = None
    sizeBytes: Optional[int] = None
    createdAt: str
    chunkCount: int

class ActiveSourcesUpdate(BaseModel):
    sourceIds: List[str]

class UrlSourceCreate(BaseModel):
    url: str

class UserPromptUpdate(BaseModel):
    customPrompt: Optional[str] = None

class UserPromptResponse(BaseModel):
    userId: str
    customPrompt: Optional[str] = None
    updatedAt: str

class ImageGenerateRequest(BaseModel):
    prompt: str
    size: Optional[str] = "1024x1024"

class GeneratedImageResponse(BaseModel):
    id: str
    projectId: str
    prompt: str
    imagePath: str
    imageUrl: str
    size: str
    createdAt: str

# ==================== HELPERS ====================

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: str, email: str) -> str:
    expiration = datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRATION_HOURS)
    payload = {
        "sub": user_id,
        "email": email,
        "exp": expiration
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

def decode_token(token: str) -> dict:
    try:
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        return payload
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token has expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    token = credentials.credentials
    payload = decode_token(token)
    user = await db.users.find_one({"id": payload["sub"]}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=401, detail="User not found")
    return user

def is_admin(email: str) -> bool:
    return email.endswith(ADMIN_EMAIL_DOMAIN)

# ==================== PERMISSION MATRIX ====================

import hashlib

def hash_string(s: str) -> str:
    """Create short hash of string for cache key"""
    return hashlib.sha256(s.encode()).hexdigest()[:16]

async def get_user_project_role(user_id: str, project: dict) -> Optional[str]:
    """Get user's role in a project. Returns None if no access."""
    if project["ownerId"] == user_id:
        return "owner"
    
    # Check sharedMembers first (new format with roles)
    shared_members = project.get("sharedMembers", [])
    for member in shared_members:
        if member.get("userId") == user_id:
            return member.get("role", "viewer")
    
    # Fallback to old sharedWith format (legacy, treat as viewer)
    if user_id in project.get("sharedWith", []):
        return "viewer"
    
    return None

async def check_project_access(user: dict, project_id: str, required_role: str = "viewer") -> dict:
    """
    Check if user has access to project with required role.
    Roles hierarchy: owner > manager > editor > viewer
    """
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    role = await get_user_project_role(user["id"], project)
    
    if role is None:
        raise HTTPException(status_code=403, detail="Access denied to this project")
    
    # Role hierarchy check
    role_levels = {"owner": 4, "manager": 3, "editor": 2, "viewer": 1}
    user_level = role_levels.get(role, 0)
    required_level = role_levels.get(required_role, 1)
    
    if user_level < required_level:
        raise HTTPException(
            status_code=403, 
            detail=f"Insufficient permissions. Required: {required_role}, Your role: {role}"
        )
    
    return {"project": project, "role": role}

async def get_user_accessible_project_ids(user_id: str) -> List[str]:
    """Get list of all project IDs user has access to"""
    # Projects where user is owner
    owned = await db.projects.find({"ownerId": user_id}, {"id": 1, "_id": 0}).to_list(1000)
    
    # Projects shared with user (old format)
    shared_old = await db.projects.find(
        {"sharedWith": user_id}, 
        {"id": 1, "_id": 0}
    ).to_list(1000)
    
    # Projects shared with user (new format)
    shared_new = await db.projects.find(
        {"sharedMembers.userId": user_id},
        {"id": 1, "_id": 0}
    ).to_list(1000)
    
    all_ids = set()
    for p in owned + shared_old + shared_new:
        all_ids.add(p["id"])
    
    return list(all_ids)

def can_manage_sources(role: str) -> bool:
    """Check if role allows source management"""
    return role in ["owner", "manager"]

def can_edit_chats(role: str) -> bool:
    """Check if role allows chat editing"""
    return role in ["owner", "manager", "editor"]

def can_manage_members(role: str) -> bool:
    """Check if role allows member management"""
    return role in ["owner", "manager"]

# ==================== SEMANTIC CACHE FUNCTIONS ====================

import numpy as np

def cosine_similarity(vec1: List[float], vec2: List[float]) -> float:
    """Calculate cosine similarity between two vectors"""
    a = np.array(vec1)
    b = np.array(vec2)
    return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b))

async def get_embedding(text: str) -> Optional[List[float]]:
    """Get embedding for text using OpenAI"""
    if not openai_client:
        return None
    try:
        response = openai_client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=text[:8000]  # Limit input length
        )
        return response.data[0].embedding
    except Exception as e:
        logger.error(f"Embedding error: {e}")
        return None

def build_cache_key_context(
    project_id: Optional[str],
    model: str,
    developer_prompt: str,
    user_prompt: Optional[str],
    source_ids: List[str]
) -> str:
    """Build deterministic cache context hash for ZERO DATA LEAKAGE"""
    components = [
        f"project:{project_id or 'none'}",
        f"model:{model}",
        f"dev_prompt:{hash_string(developer_prompt)}",
        f"user_prompt:{hash_string(user_prompt) if user_prompt else 'none'}",
        f"sources:{hash_string(','.join(sorted(source_ids)))}"
    ]
    return hash_string('|'.join(components))

async def find_cached_answer(
    question: str, 
    project_id: Optional[str],
    question_embedding: List[float],
    cache_context_hash: str,
    user_accessible_source_ids: List[str]
) -> Optional[dict]:
    """
    Find similar cached question with ZERO DATA LEAKAGE protection.
    Only returns cache hit if:
    1. Question is semantically similar
    2. Cache context (model, prompts, sources) matches
    3. User has access to ALL sources used in cached answer
    """
    
    # Build query - match project AND context hash
    query = {
        "projectId": project_id if project_id else None,
        "cacheContextHash": cache_context_hash
    }
    
    # Add TTL check
    cutoff_date = (datetime.now(timezone.utc) - timedelta(days=CACHE_TTL_DAYS)).isoformat()
    query["createdAt"] = {"$gte": cutoff_date}
    
    # Get all cache entries for this project context
    cache_entries = await db.semantic_cache.find(query, {"_id": 0}).to_list(500)
    
    if not cache_entries:
        return None
    
    best_match = None
    best_similarity = 0
    
    for entry in cache_entries:
        if not entry.get("embedding"):
            continue
        
        # SECURITY: Check user has access to ALL sources used in cached answer
        cached_source_ids = entry.get("sourceIds", [])
        if cached_source_ids:
            # User must have access to ALL sources used in the cached response
            if not all(sid in user_accessible_source_ids for sid in cached_source_ids):
                logger.info(f"Cache SKIP: User lacks access to some sources in cached entry {entry['id']}")
                continue
        
        similarity = cosine_similarity(question_embedding, entry["embedding"])
        
        if similarity > best_similarity and similarity >= CACHE_SIMILARITY_THRESHOLD:
            best_similarity = similarity
            best_match = entry
    
    if best_match:
        # Update hit count and last hit time
        await db.semantic_cache.update_one(
            {"id": best_match["id"]},
            {
                "$inc": {"hitCount": 1},
                "$set": {"lastHitAt": datetime.now(timezone.utc).isoformat()}
            }
        )
        return {
            "answer": best_match["answer"],
            "originalQuestion": best_match["question"],
            "similarity": best_similarity,
            "hitCount": best_match.get("hitCount", 0) + 1,
            "cacheId": best_match["id"],
            "sourceIds": best_match.get("sourceIds", [])
        }
    
    return None

async def save_to_cache(
    question: str,
    answer: str,
    project_id: Optional[str],
    embedding: List[float],
    user_id: str,
    cache_context_hash: str,
    source_ids: List[str],
    sources_used: Optional[List[dict]] = None
):
    """Save question-answer pair to semantic cache with full context"""
    cache_entry = {
        "id": str(uuid.uuid4()),
        "question": question,
        "answer": answer,
        "embedding": embedding,
        "projectId": project_id,
        "cacheContextHash": cache_context_hash,  # For ZERO DATA LEAKAGE
        "sourceIds": source_ids,  # Source IDs used in this response
        "sourcesUsed": sources_used,
        "createdBy": user_id,
        "createdAt": datetime.now(timezone.utc).isoformat(),
        "hitCount": 0,
        "lastHitAt": None
    }
    await db.semantic_cache.insert_one(cache_entry)
    logger.info(f"Cached answer for question: {question[:50]}... (context: {cache_context_hash})")

async def ensure_gpt_config():
    """Ensure GPT config singleton exists with strict active sources rules"""
    config = await db.gpt_config.find_one({"id": "1"}, {"_id": 0})
    default_prompt = """You are Claude, a helpful AI assistant by Anthropic. Use ONLY the active sources provided in context.

IMPORTANT RULES:
1. If no sources available - ask user to upload/activate files
2. Cite sources as [Source: name]
3. Be concise and accurate
4. Respond in the same language as the user's question
5. If the context seems incomplete or you can only find limited information, say: "I found limited information on this topic. Try activating more sources or ask a more specific question."
6. Never make up information not present in the sources"""
    
    if not config:
        config = {
            "id": "1",
            "model": "claude-sonnet-4-20250514",
            "developerPrompt": default_prompt,
            "updatedAt": datetime.now(timezone.utc).isoformat()
        }
        await db.gpt_config.insert_one(config)
    return config

async def verify_project_ownership(project_id: str, user_id: str):
    """Verify that the user owns or has access to the project"""
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Check if owner or shared with
    shared_with = project.get("sharedWith", [])
    if project["ownerId"] != user_id and user_id not in shared_with:
        raise HTTPException(status_code=403, detail="Not authorized to access this project")
    
    return project

# ==================== TEXT EXTRACTION ====================

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF file content"""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PDF extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from PDF: {str(e)}")

def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX file content"""
    try:
        doc = Document(io.BytesIO(file_content))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.error(f"DOCX extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from DOCX: {str(e)}")

def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from TXT/MD file content"""
    try:
        return file_content.decode('utf-8')
    except UnicodeDecodeError:
        try:
            return file_content.decode('latin-1')
        except Exception as e:
            logger.error(f"TXT extraction error: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to read text file: {str(e)}")

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint file content"""
    try:
        from pptx import Presentation
        from io import BytesIO
        
        prs = Presentation(BytesIO(file_content))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from PowerPoint: {str(e)}")

def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from Excel file content - optimized for AI search"""
    try:
        from openpyxl import load_workbook
        from io import BytesIO
        
        wb = load_workbook(BytesIO(file_content), data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]
            
            # Get all rows
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue
            
            # First row is usually headers
            headers = [str(cell).strip() if cell is not None else f"Column_{i}" for i, cell in enumerate(rows[0])]
            
            # Process data rows - create structured records
            for row_idx, row in enumerate(rows[1:], start=2):
                row_values = [str(cell).strip() if cell is not None else "" for cell in row]
                
                # Skip empty rows
                if not any(row_values):
                    continue
                
                # Create structured record: "Header1: Value1, Header2: Value2, ..."
                record_parts = []
                for header, value in zip(headers, row_values):
                    if value:  # Only include non-empty values
                        record_parts.append(f"{header}: {value}")
                
                if record_parts:
                    # Add row number for reference
                    record = f"Row {row_idx}: " + ", ".join(record_parts)
                    sheet_text.append(record)
            
            if len(sheet_text) > 1:
                # Add header info at the top
                sheet_text.insert(1, f"Columns: {', '.join(headers)}")
                text_parts.append("\n".join(sheet_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"XLSX extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from Excel: {str(e)}")

def extract_text_from_csv(file_content: bytes) -> str:
    """Extract text from CSV file content - optimized for large product catalogs"""
    import csv
    from io import StringIO
    
    # Increase field size limit for large cells
    csv.field_size_limit(10 * 1024 * 1024)  # 10MB per field
    
    try:
        # Try UTF-8 first, then fallback to other encodings
        text_content = None
        for encoding in ['utf-8', 'utf-8-sig', 'cp1251', 'latin-1']:
            try:
                text_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if text_content is None:
            text_content = file_content.decode('utf-8', errors='replace')
        
        # Normalize line endings and parse CSV
        text_content = text_content.replace('\r\n', '\n').replace('\r', '\n')
        
        # Parse CSV with proper newline handling
        reader = csv.reader(StringIO(text_content, newline=''))
        rows = list(reader)
        
        if not rows:
            return ""
        
        # Get headers
        headers = rows[0] if rows else []
        
        # Format as structured text for better GPT comprehension
        text_parts = []
        text_parts.append(f"[CSV Data - {len(rows)-1} records]")
        text_parts.append(f"Columns: {', '.join(headers)}")
        text_parts.append("")
        
        # Format each row as a structured record
        for i, row in enumerate(rows[1:], 1):
            if not any(cell.strip() for cell in row if cell):
                continue  # Skip empty rows
            
            record_parts = []
            for j, cell in enumerate(row):
                if cell and cell.strip():
                    header = headers[j] if j < len(headers) else f"Column{j+1}"
                    # Replace newlines in cell content with spaces
                    clean_cell = cell.strip().replace('\n', ' ').replace('\r', ' ')
                    record_parts.append(f"{header}: {clean_cell}")
            
            if record_parts:
                text_parts.append(f"[Record {i}]")
                text_parts.append(" | ".join(record_parts))
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"CSV extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from CSV: {str(e)}")

def extract_text_from_image(file_content: bytes) -> str:
    """Extract text from image using OCR (pytesseract)"""
    try:
        # Set tesseract path explicitly
        pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'
        
        # Open image with PIL
        image = Image.open(io.BytesIO(file_content))
        
        # Convert to RGB if necessary (for PNG with transparency)
        if image.mode in ('RGBA', 'P'):
            image = image.convert('RGB')
        
        # Run OCR with Russian and English languages
        text = pytesseract.image_to_string(image, lang='rus+eng')
        
        # Clean up the text
        text = text.strip()
        
        if not text:
            return "[Image: No text detected]"
        
        return f"[Image OCR Content]\n{text}"
    except Exception as e:
        logger.error(f"Image OCR error: {str(e)}")
        return f"[Image: OCR failed - {str(e)[:50]}]"

def extract_text_from_html(html_content: str) -> str:
    """Extract readable text from HTML content"""
    try:
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for element in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
            element.decompose()
        
        # Get text
        text = soup.get_text(separator='\n')
        
        # Clean up whitespace
        lines = [line.strip() for line in text.splitlines()]
        lines = [line for line in lines if line]
        
        return '\n\n'.join(lines)
    except Exception as e:
        logger.error(f"HTML extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from URL: {str(e)}")

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE) -> List[str]:
    """Split text into chunks of approximately chunk_size characters"""
    if not text:
        return []
    
    chunks = []
    current_chunk = ""
    
    # Split by paragraphs first
    paragraphs = text.split('\n\n')
    
    for para in paragraphs:
        if len(current_chunk) + len(para) + 2 <= chunk_size:
            if current_chunk:
                current_chunk += "\n\n" + para
            else:
                current_chunk = para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            
            # If paragraph is longer than chunk_size, split it
            if len(para) > chunk_size:
                words = para.split()
                current_chunk = ""
                for word in words:
                    if len(current_chunk) + len(word) + 1 <= chunk_size:
                        if current_chunk:
                            current_chunk += " " + word
                        else:
                            current_chunk = word
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = word
            else:
                current_chunk = para
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

def chunk_tabular_text(text: str, chunk_size: int = CHUNK_SIZE_TABULAR) -> List[str]:
    """Split tabular text (Excel, CSV) into chunks - preserves row integrity"""
    if not text:
        return []
    
    chunks = []
    current_chunk = ""
    
    # Split by lines (each line is a record)
    lines = text.split('\n')
    
    for line in lines:
        # Sheet headers and column info should start new chunks
        if line.startswith('[Sheet:') or line.startswith('Columns:'):
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = line
        elif len(current_chunk) + len(line) + 1 <= chunk_size:
            if current_chunk:
                current_chunk += "\n" + line
            else:
                current_chunk = line
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = line
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks

def score_chunk_relevance(chunk_content: str, query: str) -> float:
    """Score chunk relevance using simple keyword overlap"""
    # Tokenize query and chunk
    query_words = set(re.findall(r'\w+', query.lower()))
    chunk_words = set(re.findall(r'\w+', chunk_content.lower()))
    
    if not query_words:
        return 0.0
    
    # Calculate overlap
    overlap = len(query_words & chunk_words)
    return overlap / len(query_words)

async def get_relevant_chunks(source_ids: List[str], project_id: str, query: str, department_ids: List[str] = None) -> List[dict]:
    """Get most relevant chunks from active sources using keyword ranking"""
    if not source_ids:
        return []
    
    # Build projectId filter to include project, global, and department sources
    project_id_filter = [GLOBAL_PROJECT_ID]
    if project_id:
        project_id_filter.append(project_id)
    if department_ids:
        project_id_filter.extend(department_ids)
    
    logger.info(f"get_relevant_chunks: source_ids count={len(source_ids)}, project_id_filter={project_id_filter}")
    
    # Get all chunks from active sources (include project, department, and global sources)
    # Use higher limit to ensure we capture all sources including smaller ones
    all_chunks = await db.source_chunks.find({
        "sourceId": {"$in": source_ids},
        "projectId": {"$in": project_id_filter}
    }, {"_id": 0}).to_list(50000)
    
    if not all_chunks:
        return []
    
    # Score each chunk (handle both 'content' and 'text' field names)
    scored_chunks = []
    for chunk in all_chunks:
        chunk_content = chunk.get("content") or chunk.get("text", "")
        score = score_chunk_relevance(chunk_content, query)
        scored_chunks.append({**chunk, "score": score, "_content": chunk_content})
    
    # Sort by score descending
    scored_chunks.sort(key=lambda x: x["score"], reverse=True)
    
    # Select top chunks up to MAX_CHUNKS_PER_QUERY, respecting MAX_CONTEXT_CHARS
    selected_chunks = []
    total_chars = 0
    
    for chunk in scored_chunks[:MAX_CHUNKS_PER_QUERY]:  # Only consider top N chunks
        chunk_content = chunk.get("_content", "")
        if total_chars + len(chunk_content) > MAX_CONTEXT_CHARS:
            break  # Stop if we exceed character limit
        selected_chunks.append(chunk)
        total_chars += len(chunk_content)
    
    logger.info(f"RAG optimization: Selected {len(selected_chunks)}/{len(scored_chunks)} chunks, {total_chars} chars")
    
    return selected_chunks

def extract_urls_from_text(text: str) -> List[str]:
    """Extract unique URLs from text"""
    urls = URL_PATTERN.findall(text)
    # Clean URLs (remove trailing punctuation)
    cleaned = []
    for url in urls:
        # Remove trailing punctuation that might have been captured
        url = url.rstrip('.,;:!?)]}"\'')
        if url and url not in cleaned:
            cleaned.append(url)
    return cleaned[:MAX_AUTO_INGEST_URLS]  # Limit number of auto-ingested URLs

async def auto_ingest_url(url: str, project_id: str) -> Optional[dict]:
    """
    Auto-ingest a URL: fetch, extract text, chunk, and store.
    Returns the source document or None if failed.
    """
    try:
        # Check if URL already exists in this project
        existing = await db.sources.find_one({
            "url": url,
            "projectId": project_id
        })
        
        if existing:
            logger.info(f"URL already ingested: {url}")
            return existing
        
        # Fetch URL content
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True, verify=False) as http_client:
            response = await http_client.get(url, headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            })
            response.raise_for_status()
        
        # Check content type
        content_type = response.headers.get('content-type', '')
        if 'text/html' not in content_type and 'text/plain' not in content_type:
            logger.warning(f"URL {url} returned non-text content: {content_type}")
            return None
        
        # Extract text from HTML
        html_content = response.text
        extracted_text = extract_text_from_html(html_content)
        
        if not extracted_text or len(extracted_text.strip()) < 10:
            logger.warning(f"Could not extract text from URL: {url}")
            return None
        
        # Generate source ID
        source_id = str(uuid.uuid4())
        
        # Create chunks
        chunks = chunk_text(extracted_text)
        
        # Extract domain for display name
        from urllib.parse import urlparse
        parsed_url = urlparse(url)
        display_name = f"{parsed_url.netloc}{parsed_url.path[:50]}"
        
        # Save source metadata
        source_doc = {
            "id": source_id,
            "projectId": project_id,
            "kind": "url",
            "originalName": display_name,
            "url": url,
            "mimeType": "text/html",
            "sizeBytes": len(html_content.encode('utf-8')),
            "storagePath": None,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.sources.insert_one(source_doc)
        
        # Save chunks
        for i, chunk_content in enumerate(chunks):
            chunk_doc = {
                "id": str(uuid.uuid4()),
                "sourceId": source_id,
                "projectId": project_id,
                "chunkIndex": i,
                "content": chunk_content,
                "createdAt": datetime.now(timezone.utc).isoformat()
            }
            await db.source_chunks.insert_one(chunk_doc)
        
        logger.info(f"Auto-ingested URL {url} with {len(chunks)} chunks for project {project_id}")
        return source_doc
        
    except httpx.TimeoutException:
        logger.warning(f"Timeout fetching URL: {url}")
        return None
    except httpx.HTTPStatusError as e:
        logger.warning(f"HTTP error fetching URL {url}: {e.response.status_code}")
        return None
    except Exception as e:
        logger.error(f"Error auto-ingesting URL {url}: {str(e)}")
        return None

# ==================== AUTH ENDPOINTS ====================

@api_router.post("/auth/login", response_model=TokenResponse)
async def login(user_data: UserLogin):
    user = await db.users.find_one({"email": user_data.email}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=401, detail="Invalid credentials")
    
    if not verify_password(user_data.password, user["passwordHash"]):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    
    token = create_token(user["id"], user["email"])
    
    return TokenResponse(
        token=token,
        user=UserResponse(
            id=user["id"],
            email=user["email"],
            isAdmin=is_admin(user["email"]),
            createdAt=user["createdAt"]
        )
    )

@api_router.get("/auth/me", response_model=UserResponse)
async def get_me(current_user: dict = Depends(get_current_user)):
    return UserResponse(
        id=current_user["id"],
        email=current_user["email"],
        isAdmin=is_admin(current_user["email"]),
        createdAt=current_user["createdAt"],
        canEditGlobalSources=current_user.get("canEditGlobalSources", False),
        departments=current_user.get("departments", []),
        primaryDepartmentId=current_user.get("primaryDepartmentId")
    )

# ==================== PROJECT ENDPOINTS ====================

async def verify_project_access(project_id: str, user_id: str):
    """Verify user has access to project (owner or shared)"""
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    shared_with = project.get("sharedWith", [])
    if project["ownerId"] != user_id and user_id not in shared_with:
        raise HTTPException(status_code=403, detail="Not authorized to access this project")
    
    return project

@api_router.get("/projects")
async def get_projects(
    page: int = Query(1, ge=1, description="Page number"),
    page_size: int = Query(50, ge=1, le=100, description="Items per page"),
    current_user: dict = Depends(get_current_user)
):
    """Get projects with pagination"""
    query = {"$or": [{"ownerId": current_user["id"]}, {"sharedWith": current_user["id"]}]}
    result = await paginate_query(db.projects, query, page, page_size)
    # Transform items to ProjectResponse format
    result["items"] = [{**p, "sharedWith": p.get("sharedWith", [])} for p in result["items"]]
    return result

@api_router.post("/projects", response_model=ProjectResponse)
async def create_project(project_data: ProjectCreate, current_user: dict = Depends(get_current_user)):
    project_id = str(uuid.uuid4())
    project = {
        "id": project_id,
        "name": project_data.name,
        "ownerId": current_user["id"],
        "sharedWith": [],
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.projects.insert_one(project)
    return ProjectResponse(**project)

@api_router.get("/projects/{project_id}", response_model=ProjectResponse)
async def get_project(project_id: str, current_user: dict = Depends(get_current_user)):
    project = await verify_project_access(project_id, current_user["id"])
    return ProjectResponse(**{**project, "sharedWith": project.get("sharedWith", [])})

@api_router.post("/projects/{project_id}/share")
async def share_project(project_id: str, data: ShareProjectRequest, current_user: dict = Depends(get_current_user)):
    """Share project with another user by email with specified role"""
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Check if user can manage members (owner or manager)
    role = await get_user_project_role(current_user["id"], project)
    if not can_manage_members(role) and role != "owner":
        raise HTTPException(status_code=403, detail="Only owner or manager can share project")
    
    # Find user by email
    user_to_share = await db.users.find_one({"email": data.email}, {"_id": 0})
    if not user_to_share:
        raise HTTPException(status_code=404, detail="User not found")
    
    if user_to_share["id"] == current_user["id"]:
        raise HTTPException(status_code=400, detail="Cannot share with yourself")
    
    if user_to_share["id"] == project["ownerId"]:
        raise HTTPException(status_code=400, detail="Cannot change owner's role")
    
    # Validate role
    valid_roles = ["viewer", "editor", "manager"]
    share_role = data.role if data.role in valid_roles else "viewer"
    
    # Managers cannot grant manager role (only owners can)
    if share_role == "manager" and role != "owner":
        raise HTTPException(status_code=403, detail="Only owner can grant manager role")
    
    # Update sharedMembers (new format with roles)
    shared_members = project.get("sharedMembers", [])
    
    # Remove existing entry if any
    shared_members = [m for m in shared_members if m.get("userId") != user_to_share["id"]]
    
    # Add with new role
    shared_members.append({
        "userId": user_to_share["id"],
        "email": user_to_share["email"],
        "role": share_role
    })
    
    # Also update legacy sharedWith for backward compatibility
    shared_with = project.get("sharedWith", [])
    if user_to_share["id"] not in shared_with:
        shared_with.append(user_to_share["id"])
    
    await db.projects.update_one(
        {"id": project_id},
        {"$set": {"sharedWith": shared_with, "sharedMembers": shared_members}}
    )
    
    return {
        "message": f"Project shared with {data.email} as {share_role}", 
        "sharedMembers": shared_members
    }

@api_router.put("/projects/{project_id}/members/{user_id}/role")
async def update_member_role(
    project_id: str, 
    user_id: str, 
    role: str,
    current_user: dict = Depends(get_current_user)
):
    """Update a member's role in the project"""
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Only owner can change roles
    if project["ownerId"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="Only owner can change member roles")
    
    if user_id == project["ownerId"]:
        raise HTTPException(status_code=400, detail="Cannot change owner's role")
    
    valid_roles = ["viewer", "editor", "manager"]
    if role not in valid_roles:
        raise HTTPException(status_code=400, detail=f"Invalid role. Valid: {valid_roles}")
    
    # Update role in sharedMembers
    shared_members = project.get("sharedMembers", [])
    updated = False
    for member in shared_members:
        if member.get("userId") == user_id:
            member["role"] = role
            updated = True
            break
    
    if not updated:
        raise HTTPException(status_code=404, detail="Member not found in project")
    
    await db.projects.update_one(
        {"id": project_id},
        {"$set": {"sharedMembers": shared_members}}
    )
    
    return {"message": f"Role updated to {role}", "sharedMembers": shared_members}

@api_router.delete("/projects/{project_id}/share/{user_id}")
async def unshare_project(project_id: str, user_id: str, current_user: dict = Depends(get_current_user)):
    """Remove user from shared project"""
    project = await db.projects.find_one({"id": project_id}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Check if user can manage members
    role = await get_user_project_role(current_user["id"], project)
    if not can_manage_members(role) and role != "owner":
        raise HTTPException(status_code=403, detail="Only owner or manager can remove members")
    
    # Managers cannot remove other managers
    target_role = await get_user_project_role(user_id, project)
    if role == "manager" and target_role == "manager":
        raise HTTPException(status_code=403, detail="Managers cannot remove other managers")
    
    # Remove from both old and new format
    shared_with = project.get("sharedWith", [])
    if user_id in shared_with:
        shared_with.remove(user_id)
    
    shared_members = project.get("sharedMembers", [])
    shared_members = [m for m in shared_members if m.get("userId") != user_id]
    
    await db.projects.update_one(
        {"id": project_id},
        {"$set": {"sharedWith": shared_with, "sharedMembers": shared_members}}
    )
    
    return {"message": "User removed from project", "sharedMembers": shared_members}

@api_router.get("/projects/{project_id}/members")
async def get_project_members(project_id: str, current_user: dict = Depends(get_current_user)):
    """Get all members of a project with their roles"""
    project = await verify_project_access(project_id, current_user["id"])
    
    members = []
    
    # Get owner
    owner = await db.users.find_one({"id": project["ownerId"]}, {"_id": 0, "passwordHash": 0})
    if owner:
        members.append({"id": owner["id"], "email": owner["email"], "role": "owner"})
    
    # Get shared users from sharedMembers (new format with roles)
    shared_members = project.get("sharedMembers", [])
    seen_user_ids = set()
    
    for member in shared_members:
        user_id = member.get("userId")
        if user_id and user_id not in seen_user_ids:
            seen_user_ids.add(user_id)
            members.append({
                "id": user_id,
                "email": member.get("email", ""),
                "role": member.get("role", "viewer")
            })
    
    # Fallback: check sharedWith for users not in sharedMembers (legacy)
    for user_id in project.get("sharedWith", []):
        if user_id not in seen_user_ids:
            user = await db.users.find_one({"id": user_id}, {"_id": 0, "passwordHash": 0})
            if user:
                members.append({"id": user["id"], "email": user["email"], "role": "viewer"})
                seen_user_ids.add(user_id)
    
    return members

@api_router.delete("/projects/{project_id}")
async def delete_project(project_id: str, current_user: dict = Depends(get_current_user)):
    project = await db.projects.find_one({"id": project_id, "ownerId": current_user["id"]})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found or not owner")
    
    # Delete all messages in all chats of this project
    chats = await db.chats.find({"projectId": project_id}).to_list(1000)
    chat_ids = [chat["id"] for chat in chats]
    if chat_ids:
        await db.messages.delete_many({"chatId": {"$in": chat_ids}})
    
    # Delete all chats
    await db.chats.delete_many({"projectId": project_id})
    
    # Delete all sources and chunks for this project
    sources = await db.sources.find({"projectId": project_id}).to_list(1000)
    for source in sources:
        if source.get("storagePath"):
            file_path = UPLOAD_DIR / source["storagePath"]
            if file_path.exists():
                file_path.unlink()
    await db.sources.delete_many({"projectId": project_id})
    await db.source_chunks.delete_many({"projectId": project_id})
    
    # Delete project
    await db.projects.delete_one({"id": project_id})
    
    return {"message": "Project deleted successfully"}

# ==================== CHAT ENDPOINTS ====================

# --- Quick Chats (no project) ---

@api_router.get("/quick-chats")
async def get_quick_chats(
    page: int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=100),
    current_user: dict = Depends(get_current_user)
):
    """Get all quick chats with pagination"""
    query = {"ownerId": current_user["id"], "projectId": None}
    result = await paginate_query(db.chats, query, page, page_size)
    result["items"] = [{**c, "activeSourceIds": c.get("activeSourceIds", [])} for c in result["items"]]
    return result

@api_router.post("/quick-chats", response_model=ChatResponse)
async def create_quick_chat(chat_data: QuickChatCreate, current_user: dict = Depends(get_current_user)):
    """Create a quick chat without a project"""
    chat_id = str(uuid.uuid4())
    chat = {
        "id": chat_id,
        "projectId": None,
        "ownerId": current_user["id"],
        "name": chat_data.name or "Quick Chat",
        "activeSourceIds": [],
        "sourceMode": "all",
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.chats.insert_one(chat)
    return ChatResponse(**chat)

@api_router.post("/chats/{chat_id}/move", response_model=ChatResponse)
async def move_chat_to_project(chat_id: str, data: MoveChatRequest, current_user: dict = Depends(get_current_user)):
    """Move a quick chat to a project"""
    # Get the chat
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify ownership - check both ownerId (for quick chats) and projectId (for project chats)
    if chat.get("projectId"):
        await verify_project_ownership(chat["projectId"], current_user["id"])
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized")
    
    # Verify target project ownership
    await verify_project_ownership(data.targetProjectId, current_user["id"])
    
    # Update the chat
    await db.chats.update_one(
        {"id": chat_id},
        {"$set": {"projectId": data.targetProjectId, "ownerId": None}}
    )
    
    updated_chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    return ChatResponse(**{**updated_chat, "activeSourceIds": updated_chat.get("activeSourceIds", [])})

@api_router.put("/chats/{chat_id}/rename", response_model=ChatResponse)
async def rename_chat(chat_id: str, data: RenameChatRequest, current_user: dict = Depends(get_current_user)):
    """Rename a chat"""
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify ownership
    if chat.get("projectId"):
        await verify_project_ownership(chat["projectId"], current_user["id"])
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized")
    
    # Update the chat name
    await db.chats.update_one(
        {"id": chat_id},
        {"$set": {"name": data.name.strip()}}
    )
    
    updated_chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    return ChatResponse(**{**updated_chat, "activeSourceIds": updated_chat.get("activeSourceIds", [])})

# --- Project Chats ---

@api_router.get("/projects/{project_id}/chats")
async def get_chats(
    project_id: str, 
    page: int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=100),
    current_user: dict = Depends(get_current_user)
):
    """Get project chats with pagination"""
    project = await verify_project_ownership(project_id, current_user["id"])
    
    query = {"projectId": project_id}
    result = await paginate_query(db.chats, query, page, page_size)
    
    # If user is owner, return all chats
    if project["ownerId"] == current_user["id"]:
        result["items"] = [{**c, "activeSourceIds": c.get("activeSourceIds", []), "sharedWithUsers": c.get("sharedWithUsers")} for c in result["items"]]
    else:
        # If user is shared, filter chats by visibility
        visible = []
        for c in result["items"]:
            shared_with = c.get("sharedWithUsers")
            if shared_with is None or current_user["id"] in shared_with:
                visible.append({**c, "activeSourceIds": c.get("activeSourceIds", []), "sharedWithUsers": shared_with})
        result["items"] = visible
        result["total"] = len(visible)
    
    return result

@api_router.post("/projects/{project_id}/chats", response_model=ChatResponse)
async def create_chat(project_id: str, chat_data: ChatCreate, current_user: dict = Depends(get_current_user)):
    await verify_project_ownership(project_id, current_user["id"])
    
    chat_id = str(uuid.uuid4())
    chat = {
        "id": chat_id,
        "projectId": project_id,
        "name": chat_data.name or "New Chat",
        "activeSourceIds": [],
        "sourceMode": "all",
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.chats.insert_one(chat)
    return ChatResponse(**chat)

@api_router.get("/chats/{chat_id}", response_model=ChatResponse)
async def get_chat(chat_id: str, current_user: dict = Depends(get_current_user)):
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify ownership - for quick chats check ownerId, for project chats check project ownership
    if chat.get("projectId"):
        await verify_project_ownership(chat["projectId"], current_user["id"])
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized")
    
    return ChatResponse(**{**chat, "activeSourceIds": chat.get("activeSourceIds", []), "sharedWithUsers": chat.get("sharedWithUsers")})

@api_router.put("/chats/{chat_id}/visibility")
async def update_chat_visibility(chat_id: str, data: UpdateChatVisibilityRequest, current_user: dict = Depends(get_current_user)):
    """Update which shared users can see this chat (owner only)"""
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    if not chat.get("projectId"):
        raise HTTPException(status_code=400, detail="Quick chats cannot be shared")
    
    # Only project owner can change visibility
    project = await db.projects.find_one({"id": chat["projectId"]}, {"_id": 0})
    if not project or project["ownerId"] != current_user["id"]:
        raise HTTPException(status_code=403, detail="Only project owner can change chat visibility")
    
    # Update chat visibility
    # Empty list [] means hidden from all shared users
    # List with IDs means only those users can see it
    # To make visible to all, send list with all shared user IDs
    await db.chats.update_one(
        {"id": chat_id},
        {"$set": {"sharedWithUsers": data.sharedWithUsers}}
    )
    
    updated_chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    return ChatResponse(**{**updated_chat, "activeSourceIds": updated_chat.get("activeSourceIds", []), "sharedWithUsers": updated_chat.get("sharedWithUsers")})

@api_router.delete("/chats/{chat_id}")
async def delete_chat(chat_id: str, current_user: dict = Depends(get_current_user)):
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify ownership - for quick chats check ownerId, for project chats check project ownership
    if chat.get("projectId"):
        await verify_project_ownership(chat["projectId"], current_user["id"])
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized")
    
    await db.messages.delete_many({"chatId": chat_id})
    await db.chats.delete_one({"id": chat_id})
    
    return {"message": "Chat deleted successfully"}
    
    return {"message": "Chat deleted successfully"}

# ==================== SOURCE ENDPOINTS (FILES + URLS) ====================

@api_router.post("/projects/{project_id}/sources/upload", response_model=SourceResponse)
async def upload_source(
    project_id: str,
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    """Upload a file source (PDF, DOCX, TXT, MD) to a project"""
    # Permission check: require manager or owner role
    access = await check_project_access(current_user, project_id, required_role="manager")
    if not can_manage_sources(access["role"]):
        raise HTTPException(status_code=403, detail="Only owners and managers can upload sources")
    
    # Validate file type
    if file.content_type not in SUPPORTED_MIME_TYPES:
        raise HTTPException(
            status_code=400, 
            detail="Unsupported file type. Supported: PDF, DOCX, PPTX, XLSX, TXT, MD, PNG, JPEG"
        )
    
    # Read file content
    content = await file.read()
    
    # Check file size
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail=f"File size exceeds maximum of {MAX_FILE_SIZE // (1024*1024)}MB")
    
    # Extract text based on file type
    file_type = SUPPORTED_MIME_TYPES[file.content_type]
    
    if file_type == "pdf":
        extracted_text = extract_text_from_pdf(content)
    elif file_type == "docx":
        extracted_text = extract_text_from_docx(content)
    elif file_type == "pptx":
        extracted_text = extract_text_from_pptx(content)
    elif file_type == "xlsx":
        extracted_text = extract_text_from_xlsx(content)
    elif file_type == "csv":
        extracted_text = extract_text_from_csv(content)
    elif file_type in ["png", "jpeg", "jpg"]:
        # Use OCR to extract text from images
        extracted_text = extract_text_from_image(content)
    else:  # txt or md
        extracted_text = extract_text_from_txt(content)
    
    # For non-image files, check if text was extracted
    if file_type not in ["png", "jpeg", "jpg"]:
        if not extracted_text or len(extracted_text.strip()) < 10:
            raise HTTPException(
                status_code=400, 
                detail="This file appears to be empty or contains no extractable text. For PDFs, please ensure it's text-based, not image-based/scanned."
            )
    
    # Generate source ID and storage path
    source_id = str(uuid.uuid4())
    storage_filename = f"{source_id}.{file_type}"
    storage_path = UPLOAD_DIR / storage_filename
    
    # Save file to disk
    async with aiofiles.open(storage_path, 'wb') as f:
        await f.write(content)
    
    # Create chunks - use tabular chunking for Excel/CSV
    if file_type in ["xlsx", "csv"]:
        chunks = chunk_tabular_text(extracted_text)
    else:
        chunks = chunk_text(extracted_text)
    
    # Save source metadata
    source_doc = {
        "id": source_id,
        "projectId": project_id,
        "kind": "file",
        "originalName": file.filename,
        "url": None,
        "mimeType": file.content_type,
        "sizeBytes": len(content),
        "storagePath": storage_filename,
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.sources.insert_one(source_doc)
    
    # Save chunks
    for i, chunk_content in enumerate(chunks):
        chunk_doc = {
            "id": str(uuid.uuid4()),
            "sourceId": source_id,
            "projectId": project_id,
            "chunkIndex": i,
            "content": chunk_content,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.source_chunks.insert_one(chunk_doc)
    
    logger.info(f"Uploaded file {file.filename} with {len(chunks)} chunks for project {project_id}")
    
    return SourceResponse(
        id=source_id,
        projectId=project_id,
        kind="file",
        originalName=file.filename,
        url=None,
        mimeType=file.content_type,
        sizeBytes=len(content),
        createdAt=source_doc["createdAt"],
        chunkCount=len(chunks)
    )

@api_router.post("/projects/{project_id}/sources/upload-multiple")
async def upload_multiple_sources(
    project_id: str,
    files: List[UploadFile] = File(...),
    current_user: dict = Depends(get_current_user)
):
    """Upload multiple file sources to a project"""
    await verify_project_access(project_id, current_user["id"])
    
    results = []
    errors = []
    
    for file in files:
        try:
            # Validate file type
            if file.content_type not in SUPPORTED_MIME_TYPES:
                errors.append({"filename": file.filename, "error": "Unsupported file type"})
                continue
            
            # Read file content
            content = await file.read()
            
            # Check file size
            if len(content) > MAX_FILE_SIZE:
                errors.append({"filename": file.filename, "error": f"File size exceeds maximum"})
                continue
            
            # Extract text based on file type
            file_type = SUPPORTED_MIME_TYPES[file.content_type]
            
            try:
                if file_type == "pdf":
                    extracted_text = extract_text_from_pdf(content)
                elif file_type == "docx":
                    extracted_text = extract_text_from_docx(content)
                elif file_type == "pptx":
                    extracted_text = extract_text_from_pptx(content)
                elif file_type == "xlsx":
                    extracted_text = extract_text_from_xlsx(content)
                elif file_type == "csv":
                    extracted_text = extract_text_from_csv(content)
                elif file_type in ["png", "jpeg", "jpg"]:
                    extracted_text = extract_text_from_image(content)
                else:
                    extracted_text = extract_text_from_txt(content)
            except Exception as e:
                errors.append({"filename": file.filename, "error": str(e)[:100]})
                continue
            
            # For non-image files, check if text was extracted
            if file_type not in ["png", "jpeg", "jpg"]:
                if not extracted_text or len(extracted_text.strip()) < 10:
                    errors.append({"filename": file.filename, "error": "No extractable text"})
                    continue
            
            # Generate source ID and storage path
            source_id = str(uuid.uuid4())
            storage_filename = f"{source_id}.{file_type}"
            storage_path = UPLOAD_DIR / storage_filename
            
            # Save file to disk
            async with aiofiles.open(storage_path, 'wb') as f:
                await f.write(content)
            
            # Create chunks - use tabular chunking for Excel/CSV
            if file_type in ["xlsx", "csv"]:
                chunks = chunk_tabular_text(extracted_text)
            else:
                chunks = chunk_text(extracted_text)
            
            # Save source metadata
            source_doc = {
                "id": source_id,
                "projectId": project_id,
                "kind": "file",
                "originalName": file.filename,
                "url": None,
                "mimeType": file.content_type,
                "sizeBytes": len(content),
                "storagePath": storage_filename,
                "createdAt": datetime.now(timezone.utc).isoformat()
            }
            await db.sources.insert_one(source_doc)
            
            # Save chunks
            for i, chunk_content in enumerate(chunks):
                chunk_doc = {
                    "id": str(uuid.uuid4()),
                    "sourceId": source_id,
                    "projectId": project_id,
                    "chunkIndex": i,
                    "content": chunk_content,
                    "createdAt": datetime.now(timezone.utc).isoformat()
                }
                await db.source_chunks.insert_one(chunk_doc)
            
            results.append({
                "id": source_id,
                "filename": file.filename,
                "chunkCount": len(chunks)
            })
            
        except Exception as e:
            errors.append({"filename": file.filename, "error": str(e)[:100]})
    
    return {"uploaded": results, "errors": errors}

@api_router.post("/projects/{project_id}/sources/url", response_model=SourceResponse)
async def add_url_source(
    project_id: str,
    data: UrlSourceCreate,
    current_user: dict = Depends(get_current_user)
):
    """Add a URL source to a project"""
    await verify_project_ownership(project_id, current_user["id"])
    
    url = data.url.strip()
    
    # Validate URL format
    if not url.startswith(('http://', 'https://')):
        raise HTTPException(status_code=400, detail="URL must start with http:// or https://")
    
    # Fetch URL content
    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True, verify=False) as http_client:
            response = await http_client.get(url, headers={
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            })
            response.raise_for_status()
    except httpx.TimeoutException:
        raise HTTPException(status_code=400, detail="URL fetch timed out")
    except httpx.HTTPStatusError as e:
        raise HTTPException(status_code=400, detail=f"URL returned error: {e.response.status_code}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {str(e)[:100]}")
    
    # Check content type
    content_type = response.headers.get('content-type', '')
    if 'text/html' not in content_type and 'text/plain' not in content_type:
        raise HTTPException(status_code=400, detail="URL must return HTML or text content")
    
    # Extract text from HTML
    html_content = response.text
    extracted_text = extract_text_from_html(html_content)
    
    if not extracted_text or len(extracted_text.strip()) < 10:
        raise HTTPException(status_code=400, detail="Could not extract meaningful text from this URL")
    
    # Generate source ID
    source_id = str(uuid.uuid4())
    
    # Create chunks
    chunks = chunk_text(extracted_text)
    
    # Extract domain for display name
    from urllib.parse import urlparse
    parsed_url = urlparse(url)
    display_name = f"{parsed_url.netloc}{parsed_url.path[:50]}"
    
    # Save source metadata
    source_doc = {
        "id": source_id,
        "projectId": project_id,
        "kind": "url",
        "originalName": display_name,
        "url": url,
        "mimeType": "text/html",
        "sizeBytes": len(html_content.encode('utf-8')),
        "storagePath": None,
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.sources.insert_one(source_doc)
    
    # Save chunks
    for i, chunk_content in enumerate(chunks):
        chunk_doc = {
            "id": str(uuid.uuid4()),
            "sourceId": source_id,
            "projectId": project_id,
            "chunkIndex": i,
            "content": chunk_content,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.source_chunks.insert_one(chunk_doc)
    
    logger.info(f"Added URL {url} with {len(chunks)} chunks for project {project_id}")
    
    return SourceResponse(
        id=source_id,
        projectId=project_id,
        kind="url",
        originalName=display_name,
        url=url,
        mimeType="text/html",
        sizeBytes=source_doc["sizeBytes"],
        createdAt=source_doc["createdAt"],
        chunkCount=len(chunks)
    )

@api_router.get("/projects/{project_id}/sources", response_model=List[SourceResponse])
async def list_sources(project_id: str, current_user: dict = Depends(get_current_user)):
    """List all sources in a project"""
    await verify_project_ownership(project_id, current_user["id"])
    
    sources = await db.sources.find({"projectId": project_id}, {"_id": 0}).to_list(1000)
    
    result = []
    for s in sources:
        chunk_count = await db.source_chunks.count_documents({"sourceId": s["id"]})
        result.append(SourceResponse(
            id=s["id"],
            projectId=s["projectId"],
            kind=s["kind"],
            originalName=s.get("originalName"),
            url=s.get("url"),
            mimeType=s.get("mimeType"),
            sizeBytes=s.get("sizeBytes"),
            createdAt=s["createdAt"],
            chunkCount=chunk_count
        ))
    
    return result

@api_router.delete("/projects/{project_id}/sources/{source_id}")
async def delete_source(project_id: str, source_id: str, current_user: dict = Depends(get_current_user)):
    """Delete a source from a project"""
    await verify_project_ownership(project_id, current_user["id"])
    
    source = await db.sources.find_one({"id": source_id, "projectId": project_id})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    
    # Delete physical file if exists
    if source.get("storagePath"):
        file_path = UPLOAD_DIR / source["storagePath"]
        if file_path.exists():
            file_path.unlink()
    
    # Delete chunks
    await db.source_chunks.delete_many({"sourceId": source_id})
    
    # Delete source metadata
    await db.sources.delete_one({"id": source_id})
    
    # Remove from active sources in all chats
    await db.chats.update_many(
        {"projectId": project_id},
        {"$pull": {"activeSourceIds": source_id}}
    )
    
    return {"message": "Source deleted successfully"}

class SearchRequest(BaseModel):
    query: str
    limit: Optional[int] = 20

class SearchResult(BaseModel):
    sourceId: str
    sourceName: str
    sourceKind: str
    chunkIndex: int
    content: str
    matchCount: int

@api_router.post("/projects/{project_id}/sources/search")
async def search_sources(project_id: str, search_data: SearchRequest, current_user: dict = Depends(get_current_user)):
    """Search through all source chunks in project - NO GPT, NO TOKENS"""
    await verify_project_access(project_id, current_user["id"])
    
    query = search_data.query.strip().lower()
    if not query or len(query) < 2:
        raise HTTPException(status_code=400, detail="Search query must be at least 2 characters")
    
    # Get all sources in project
    sources = await db.sources.find({"projectId": project_id}, {"_id": 0}).to_list(1000)
    source_map = {s["id"]: s for s in sources}
    
    # Search through chunks (handle both 'content' and 'text' field names)
    results = []
    chunks = await db.source_chunks.find(
        {"projectId": project_id},
        {"_id": 0}
    ).to_list(10000)
    
    for chunk in chunks:
        chunk_content = chunk.get("content") or chunk.get("text", "")
        content_lower = chunk_content.lower()
        if query in content_lower:
            source = source_map.get(chunk["sourceId"], {})
            
            # Count matches
            match_count = content_lower.count(query)
            
            # Get snippet around first match (150 chars before and after)
            idx = content_lower.find(query)
            start = max(0, idx - 150)
            end = min(len(chunk_content), idx + len(query) + 150)
            snippet = chunk_content[start:end]
            if start > 0:
                snippet = "..." + snippet
            if end < len(chunk_content):
                snippet = snippet + "..."
            
            results.append({
                "sourceId": chunk["sourceId"],
                "sourceName": source.get("originalName") or source.get("url", "Unknown"),
                "sourceKind": source.get("kind", "file"),
                "chunkIndex": chunk.get("chunkIndex", 0),
                "content": snippet,
                "matchCount": match_count
            })
    
    # Sort by match count and limit
    results.sort(key=lambda x: x["matchCount"], reverse=True)
    return results[:search_data.limit]

@api_router.get("/projects/{project_id}/sources/{source_id}/download")
async def download_source(project_id: str, source_id: str, current_user: dict = Depends(get_current_user)):
    """Download a single source file"""
    await verify_project_access(project_id, current_user["id"])
    
    source = await db.sources.find_one({"id": source_id, "projectId": project_id}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    
    if source["kind"] != "file" or not source.get("storagePath"):
        raise HTTPException(status_code=400, detail="This source cannot be downloaded")
    
    file_path = UPLOAD_DIR / source["storagePath"]
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found on disk")
    
    return FileResponse(
        path=file_path,
        filename=source.get("originalName", source["storagePath"]),
        media_type=source.get("mimeType", "application/octet-stream")
    )

@api_router.get("/projects/{project_id}/sources/{source_id}/preview")
async def preview_source(project_id: str, source_id: str, current_user: dict = Depends(get_current_user)):
    """Get source preview - extracted text content with quality info"""
    await verify_project_access(project_id, current_user["id"])
    
    source = await db.sources.find_one({"id": source_id, "projectId": project_id}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    
    # Get chunks for this source
    chunks = await db.source_chunks.find(
        {"sourceId": source_id}, 
        {"_id": 0, "content": 1, "chunkIndex": 1}
    ).sort("chunkIndex", 1).to_list(1000)
    
    # Combine chunks into full text (handle both 'content' and 'text' field names)
    full_text = "\n\n".join([c.get("content") or c.get("text", "") for c in chunks])
    
    # Calculate quality metrics
    char_count = len(full_text)
    word_count = len(full_text.split())
    
    # Determine quality based on content
    is_image = source.get("mimeType", "").startswith("image/")
    quality = "good"
    quality_message = "Текст извлечён успешно"
    
    if char_count == 0:
        quality = "empty"
        quality_message = "Текст не извлечён"
    elif is_image:
        if "[Image: No text detected]" in full_text or "[Image: OCR failed" in full_text:
            quality = "poor"
            quality_message = "OCR не смог распознать текст. Попробуйте загрузить изображение лучшего качества"
        elif word_count < 10:
            quality = "low"
            quality_message = "Мало текста распознано. Возможно изображение низкого качества"
        else:
            quality_message = f"OCR распознал {word_count} слов"
    elif word_count < 20:
        quality = "low"
        quality_message = "Мало текста извлечено"
    
    return {
        "id": source_id,
        "name": source.get("originalName") or source.get("url"),
        "kind": source["kind"],
        "mimeType": source.get("mimeType"),
        "text": full_text,
        "chunkCount": len(chunks),
        "charCount": char_count,
        "wordCount": word_count,
        "quality": quality,
        "qualityMessage": quality_message
    }

@api_router.get("/projects/{project_id}/sources/download-all")
async def download_all_sources(project_id: str, current_user: dict = Depends(get_current_user)):
    """Download all file sources as a ZIP archive"""
    import zipfile
    from io import BytesIO
    
    await verify_project_access(project_id, current_user["id"])
    
    # Get all file sources (not URLs)
    sources = await db.sources.find({
        "projectId": project_id,
        "kind": "file"
    }, {"_id": 0}).to_list(1000)
    
    if not sources:
        raise HTTPException(status_code=404, detail="No files to download")
    
    # Create ZIP in memory
    zip_buffer = BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for source in sources:
            if source.get("storagePath"):
                file_path = UPLOAD_DIR / source["storagePath"]
                if file_path.exists():
                    # Use original filename in zip
                    filename = source.get("originalName", source["storagePath"])
                    with open(file_path, 'rb') as f:
                        zip_file.writestr(filename, f.read())
    
    zip_buffer.seek(0)
    
    return StreamingResponse(
        zip_buffer,
        media_type="application/zip",
        headers={
            "Content-Disposition": f"attachment; filename=project_files.zip"
        }
    )

# ==================== IMAGE GENERATION ENDPOINTS ====================

async def check_image_rate_limit(user_id: str) -> bool:
    """Check if user has exceeded image generation rate limit"""
    one_hour_ago = (datetime.now(timezone.utc) - timedelta(hours=1)).isoformat()
    
    count = await db.generated_images.count_documents({
        "userId": user_id,
        "createdAt": {"$gte": one_hour_ago}
    })
    
    return count < IMAGE_RATE_LIMIT_PER_HOUR

@api_router.post("/projects/{project_id}/generate-image", response_model=GeneratedImageResponse)
async def generate_image(
    project_id: str,
    request: ImageGenerateRequest,
    current_user: dict = Depends(get_current_user)
):
    """Generate an AI image using OpenAI's gpt-image-1 model"""
    await verify_project_ownership(project_id, current_user["id"])
    
    # Validate prompt
    if not request.prompt or len(request.prompt.strip()) < 3:
        raise HTTPException(status_code=400, detail="Prompt must be at least 3 characters")
    
    if len(request.prompt) > 4000:
        raise HTTPException(status_code=400, detail="Prompt must be less than 4000 characters")
    
    # Validate size
    size = request.size or "1024x1024"
    if size not in VALID_IMAGE_SIZES:
        raise HTTPException(status_code=400, detail=f"Invalid size. Valid options: {VALID_IMAGE_SIZES}")
    
    # Check rate limit
    if not await check_image_rate_limit(current_user["id"]):
        raise HTTPException(
            status_code=429, 
            detail=f"Rate limit exceeded. Maximum {IMAGE_RATE_LIMIT_PER_HOUR} images per hour."
        )
    
    # Check OpenAI client
    if not openai_client:
        raise HTTPException(status_code=500, detail="OpenAI API key not configured")
    
    try:
        logger.info(f"Generating image for project {project_id}: {request.prompt[:50]}...")
        
        # Call OpenAI Images API
        response = openai_client.images.generate(
            model="gpt-image-1",
            prompt=request.prompt,
            size=size,
            n=1
        )
        
        # Get base64 image data
        image_data = response.data[0]
        
        # Handle both b64_json and url response formats
        if hasattr(image_data, 'b64_json') and image_data.b64_json:
            import base64
            image_bytes = base64.b64decode(image_data.b64_json)
        elif hasattr(image_data, 'url') and image_data.url:
            # Fetch image from URL
            async with httpx.AsyncClient(timeout=60.0, verify=False) as http_client:
                img_response = await http_client.get(image_data.url)
                image_bytes = img_response.content
        else:
            raise HTTPException(status_code=500, detail="No image data returned from OpenAI")
        
        # Generate unique filename
        image_id = str(uuid.uuid4())
        image_filename = f"{image_id}.png"
        image_path = IMAGES_DIR / image_filename
        
        # Save image to disk
        async with aiofiles.open(image_path, 'wb') as f:
            await f.write(image_bytes)
        
        # Save metadata to DB
        image_doc = {
            "id": image_id,
            "projectId": project_id,
            "userId": current_user["id"],
            "prompt": request.prompt,
            "imagePath": image_filename,
            "size": size,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.generated_images.insert_one(image_doc)
        
        logger.info(f"Generated image {image_id} for project {project_id}")
        
        return GeneratedImageResponse(
            id=image_id,
            projectId=project_id,
            prompt=request.prompt,
            imagePath=image_filename,
            imageUrl=f"/api/images/{image_id}",
            size=size,
            createdAt=image_doc["createdAt"]
        )
        
    except Exception as e:
        logger.error(f"Image generation error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate image: {str(e)[:100]}")

@api_router.get("/images/{image_id}")
async def get_image(image_id: str, current_user: dict = Depends(get_current_user)):
    """Get a generated image by ID"""
    image = await db.generated_images.find_one({"id": image_id}, {"_id": 0})
    if not image:
        raise HTTPException(status_code=404, detail="Image not found")
    
    # Verify user owns the project
    await verify_project_ownership(image["projectId"], current_user["id"])
    
    image_path = IMAGES_DIR / image["imagePath"]
    if not image_path.exists():
        raise HTTPException(status_code=404, detail="Image file not found")
    
    from fastapi.responses import FileResponse
    return FileResponse(
        path=image_path,
        media_type="image/png",
        filename=f"generated_{image_id}.png"
    )

@api_router.get("/projects/{project_id}/images", response_model=List[GeneratedImageResponse])
async def list_project_images(project_id: str, current_user: dict = Depends(get_current_user)):
    """List all generated images in a project"""
    await verify_project_ownership(project_id, current_user["id"])
    
    images = await db.generated_images.find(
        {"projectId": project_id},
        {"_id": 0}
    ).sort("createdAt", -1).to_list(100)
    
    return [
        GeneratedImageResponse(
            id=img["id"],
            projectId=img["projectId"],
            prompt=img["prompt"],
            imagePath=img["imagePath"],
            imageUrl=f"/api/images/{img['id']}",
            size=img.get("size", "1024x1024"),
            createdAt=img["createdAt"]
        )
        for img in images
    ]

@api_router.delete("/projects/{project_id}/images/{image_id}")
async def delete_image(project_id: str, image_id: str, current_user: dict = Depends(get_current_user)):
    """Delete a generated image"""
    await verify_project_ownership(project_id, current_user["id"])
    
    image = await db.generated_images.find_one({"id": image_id, "projectId": project_id})
    if not image:
        raise HTTPException(status_code=404, detail="Image not found")
    
    # Delete physical file
    image_path = IMAGES_DIR / image["imagePath"]
    if image_path.exists():
        image_path.unlink()
    
    # Delete metadata
    await db.generated_images.delete_one({"id": image_id})
    
    return {"message": "Image deleted successfully"}

# ==================== ACTIVE SOURCES ENDPOINTS ====================

@api_router.post("/chats/{chat_id}/active-sources")
async def set_active_sources(chat_id: str, data: ActiveSourcesUpdate, current_user: dict = Depends(get_current_user)):
    """Set the active sources for a chat"""
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    await verify_project_ownership(chat["projectId"], current_user["id"])
    
    # Verify all source IDs belong to this project
    if data.sourceIds:
        sources = await db.sources.find({
            "id": {"$in": data.sourceIds},
            "projectId": chat["projectId"]
        }).to_list(1000)
        
        valid_ids = {s["id"] for s in sources}
        invalid_ids = set(data.sourceIds) - valid_ids
        
        if invalid_ids:
            raise HTTPException(status_code=400, detail=f"Invalid source IDs: {invalid_ids}")
    
    # Update chat with active source IDs
    await db.chats.update_one(
        {"id": chat_id},
        {"$set": {"activeSourceIds": data.sourceIds}}
    )
    
    return {"message": "Active sources updated", "activeSourceIds": data.sourceIds}

class SourceModeUpdate(BaseModel):
    sourceMode: str  # 'all' or 'my'

@api_router.put("/chats/{chat_id}/source-mode")
async def update_source_mode(chat_id: str, data: SourceModeUpdate, current_user: dict = Depends(get_current_user)):
    """Update source mode for a chat"""
    if data.sourceMode not in ['all', 'my']:
        raise HTTPException(status_code=400, detail="Invalid source mode. Use 'all' or 'my'")
    
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify ownership
    if chat.get("projectId"):
        await verify_project_ownership(chat["projectId"], current_user["id"])
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized")
    
    await db.chats.update_one(
        {"id": chat_id},
        {"$set": {"sourceMode": data.sourceMode}}
    )
    
    return {"message": "Source mode updated", "sourceMode": data.sourceMode}

@api_router.get("/chats/{chat_id}/active-sources")
async def get_active_sources(chat_id: str, current_user: dict = Depends(get_current_user)):
    """Get the active sources for a chat"""
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    await verify_project_ownership(chat["projectId"], current_user["id"])
    
    active_source_ids = chat.get("activeSourceIds", [])
    
    if not active_source_ids:
        return {"activeSources": []}
    
    sources = await db.sources.find({
        "id": {"$in": active_source_ids},
        "projectId": chat["projectId"]
    }, {"_id": 0}).to_list(1000)
    
    return {"activeSources": sources}

# ==================== MESSAGE ENDPOINTS ====================

@api_router.get("/chats/{chat_id}/messages", response_model=List[MessageResponse])
async def get_messages(chat_id: str, current_user: dict = Depends(get_current_user)):
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify ownership - for quick chats check ownerId, for project chats check project ownership
    if chat.get("projectId"):
        await verify_project_ownership(chat["projectId"], current_user["id"])
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized")
    
    messages = await db.messages.find({"chatId": chat_id}, {"_id": 0}).sort("createdAt", 1).to_list(1000)
    
    # Get sender info for user messages
    result = []
    for m in messages:
        msg_data = {
            **m, 
            "citations": m.get("citations"), 
            "usedSources": m.get("usedSources"),
            "autoIngestedUrls": m.get("autoIngestedUrls"),
            "senderEmail": m.get("senderEmail"),
            "senderName": m.get("senderName")
        }
        result.append(MessageResponse(**msg_data))
    
    return result

@api_router.post("/chats/{chat_id}/messages", response_model=MessageResponse)
async def send_message(chat_id: str, message_data: MessageCreate, current_user: dict = Depends(get_current_user)):
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    # Verify access - for quick chats check ownerId, for project chats use permission matrix
    project_id = chat.get("projectId")
    user_role = None
    
    if project_id:
        # Use permission matrix for project chats
        try:
            access = await check_project_access(current_user, project_id, required_role="viewer")
            user_role = access["role"]
        except HTTPException:
            raise HTTPException(status_code=403, detail="Not authorized to access this project")
    elif chat.get("ownerId") != current_user["id"]:
        raise HTTPException(status_code=403, detail="Not authorized to access this chat")
    
    # === AUTO-INGEST URLs from message (only for project chats with editor+ permission) ===
    detected_urls = extract_urls_from_text(message_data.content)
    auto_ingested_sources = []
    auto_ingest_notes = []
    
    if detected_urls and project_id and can_edit_chats(user_role):
        logger.info(f"Detected {len(detected_urls)} URL(s) in message: {detected_urls}")
        
        for url in detected_urls:
            source = await auto_ingest_url(url, project_id)
            if source:
                auto_ingested_sources.append(source)
                auto_ingest_notes.append(f"Auto-ingested: {source.get('originalName', url)}")
            else:
                auto_ingest_notes.append(f"Could not fetch: {url}")
    
    # === GET SOURCE IDS WITH HIERARCHY (Project > Department > Global) ===
    project_source_ids = []
    department_source_ids = []
    global_source_ids = []
    personal_source_ids = []
    user_department_ids = current_user.get("departments", [])  # Always get this for chunk lookup
    
    # Get source mode from chat (default to 'all')
    source_mode = chat.get("sourceMode", "all")
    
    # Personal sources (always available)
    personal_sources = await db.sources.find({
        "level": "personal",
        "ownerId": current_user["id"],
        "status": {"$in": ["active", None]}
    }, {"_id": 0, "id": 1}).to_list(1000)
    personal_source_ids = [s["id"] for s in personal_sources]
    
    # Project sources (if in a project)
    if project_id:
        project_sources = await db.sources.find({
            "projectId": project_id,
            "level": {"$in": ["project", None]},  # Legacy sources have no level
            "status": {"$in": ["active", None]}
        }, {"_id": 0, "id": 1}).to_list(1000)
        project_source_ids = [s["id"] for s in project_sources]
    
    # Only include department and global sources if source_mode is 'all'
    if source_mode == 'all':
        # Department sources (from user's departments)
        if user_department_ids:
            department_sources = await db.sources.find({
                "departmentId": {"$in": user_department_ids},
                "level": "department",
                "status": "active"
            }, {"_id": 0, "id": 1}).to_list(1000)
            department_source_ids = [s["id"] for s in department_sources]
        
        # Global sources (always included in 'all' mode)
        global_sources = await db.sources.find({
            "$or": [
                {"projectId": GLOBAL_PROJECT_ID},
                {"level": "global", "status": "active"}
            ]
        }, {"_id": 0, "id": 1}).to_list(1000)
        global_source_ids = [s["id"] for s in global_sources]
    
    # Combined list with priority order: Personal > Project > Department > Global
    active_source_ids = personal_source_ids + project_source_ids + department_source_ids + global_source_ids
    
    # Get user's accessible source IDs for cache security
    user_accessible_source_ids = personal_source_ids + project_source_ids + department_source_ids + global_source_ids
    
    # Get sender display name (use part before @ in email)
    sender_email = current_user["email"]
    sender_name = sender_email.split("@")[0] if sender_email else "User"
    
    # Save user message with sender info
    user_msg_id = str(uuid.uuid4())
    user_message = {
        "id": user_msg_id,
        "chatId": chat_id,
        "role": "user",
        "content": message_data.content,
        "citations": None,
        "autoIngestedUrls": [s["id"] for s in auto_ingested_sources] if auto_ingested_sources else None,
        "senderEmail": sender_email,
        "senderName": sender_name,
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.messages.insert_one(user_message)
    
    # Get GPT config
    config = await ensure_gpt_config()
    
    # Get chat history
    history = await db.messages.find({"chatId": chat_id}, {"_id": 0}).sort("createdAt", 1).to_list(1000)
    
    # Get active sources and relevant chunks (including newly ingested)
    citations = []
    document_context = ""
    active_source_names = []
    source_types = {}  # Track source type (project/department/global)
    
    if active_source_ids:
        # Get source info with type - need to query all sources not just by projectId
        sources = await db.sources.find({
            "id": {"$in": active_source_ids}
        }, {"_id": 0}).to_list(1000)
        
        source_names = {}
        for s in sources:
            name = s.get("originalName") or s.get("url") or "Unknown"
            source_names[s["id"]] = name
            active_source_names.append(name)
            # Track source level: project > department > global
            level = s.get("level")
            if level == "department":
                source_types[s["id"]] = "department"
            elif s.get("projectId") == GLOBAL_PROJECT_ID or level == "global":
                source_types[s["id"]] = "global"
            else:
                source_types[s["id"]] = "project"
        
        # Get relevant chunks using keyword ranking
        # Priority: project chunks first, then department, then global
        relevant_chunks = await get_relevant_chunks(
            active_source_ids, 
            project_id, 
            message_data.content,
            user_department_ids  # Pass department IDs for chunk lookup
        )
        
        if relevant_chunks:
            # Sort by priority (project > department > global) and then by score
            def chunk_priority(chunk):
                source_type = source_types.get(chunk["sourceId"], "global")
                # Priority: project=0, department=1, global=2
                type_priority = {"project": 0, "department": 1, "global": 2}.get(source_type, 2)
                return (type_priority, -chunk.get("score", 0))
            
            relevant_chunks.sort(key=chunk_priority)
            
            # Build context with chunk markers and conflict detection
            context_parts = []
            seen_topics = {}  # For conflict detection
            
            for chunk in relevant_chunks:
                source_name = source_names.get(chunk["sourceId"], "Unknown")
                source_type = source_types.get(chunk["sourceId"], "global")
                chunk_content = chunk.get("content") or chunk.get("text") or chunk.get("_content", "")
                chunk_marker = f"[Source: {source_name} ({source_type.upper()}), Chunk {chunk['chunkIndex']+1}]"
                context_parts.append(f"{chunk_marker}\n{chunk_content}")
                
                # Enhanced citation with full context
                citations.append({
                    "sourceName": source_name,
                    "sourceId": chunk["sourceId"],
                    "sourceType": source_type,  # "project" or "global"
                    "chunkId": chunk.get("id", ""),
                    "chunkIndex": chunk["chunkIndex"],
                    "textFragment": chunk_content[:200] + "..." if len(chunk_content) > 200 else chunk_content,
                    "score": chunk.get("score", 0)
                })
            
            document_context = "\n\n---\n\n".join(context_parts)
    
    # Get user's custom prompt
    user_prompt_doc = await db.user_prompts.find_one({"userId": current_user["id"]}, {"_id": 0})
    user_custom_prompt = user_prompt_doc.get("customPrompt") if user_prompt_doc else None
    
    # === BUILD CACHE CONTEXT HASH (ZERO DATA LEAKAGE) ===
    config = await ensure_gpt_config()
    user_model = current_user.get("gptModel")
    model_to_use = user_model if user_model else config["model"]
    
    cache_context_hash = build_cache_key_context(
        project_id=project_id,
        model=model_to_use,
        developer_prompt=config["developerPrompt"],
        user_prompt=user_custom_prompt,
        source_ids=active_source_ids
    )
    
    # === SEMANTIC CACHE CHECK ===
    cache_hit = None
    question_embedding = None
    cache_info = None
    
    # Only use cache if there are active sources (context-dependent answers)
    if active_source_ids and openai_client:
        question_embedding = await get_embedding(message_data.content)
        if question_embedding:
            cache_hit = await find_cached_answer(
                message_data.content,
                project_id,
                question_embedding,
                cache_context_hash,
                user_accessible_source_ids
            )
            
            if cache_hit:
                logger.info(f"Cache HIT! Similarity: {cache_hit['similarity']:.3f}, Hits: {cache_hit['hitCount']}")
                cache_info = {
                    "similarity": cache_hit['similarity'],
                    "hitCount": cache_hit['hitCount'],
                    "cacheId": cache_hit['cacheId']
                }
    
    # Prepare messages for Claude
    try:
        import anthropic
        claude_client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)
        
        from_cache = False
        
        # If cache hit, use cached answer
        if cache_hit:
            response_text = cache_hit["answer"]
            response_text += f"\n\n---\n_📦 Ответ из кэша (схожесть: {cache_hit['similarity']:.0%})_"
            tokens_used = 0
            from_cache = True
        else:
            # Build system prompt for Claude
            system_parts = [config["developerPrompt"]]
            
            # Add user's custom prompt if exists
            if user_custom_prompt:
                logger.info(f"Adding user custom prompt for user {current_user['id']}")
                system_parts.append(f"USER INSTRUCTIONS:\n{user_custom_prompt}")
            
            # Add document context if available
            if document_context:
                active_sources_list = ", ".join(active_source_names) if active_source_names else "None"
                chunks_count = len(citations) if citations else 0
                context_message = f"SOURCES: {active_sources_list}\nCHUNKS: {chunks_count} (top {MAX_CHUNKS_PER_QUERY} most relevant)\n\n{document_context[:10000]}"
                system_parts.append(context_message)
            
            system_prompt = "\n\n".join(system_parts)
            
            # Build messages for Claude (no system role in messages)
            messages = []
            for msg in history[:-1]:
                messages.append({"role": msg["role"], "content": msg["content"]})
            messages.append({"role": "user", "content": message_data.content})
            
            # Call Claude API
            claude_response = claude_client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=4096,
                system=system_prompt,
                messages=messages
            )
            
            response_text = claude_response.content[0].text
            
            # Track tokens
            tokens_used = claude_response.usage.input_tokens + claude_response.usage.output_tokens
            
            # Update user token usage in DB
            if tokens_used > 0:
                await db.token_usage.update_one(
                    {"userId": current_user["id"]},
                    {
                        "$inc": {"totalTokens": tokens_used, "messageCount": 1},
                        "$set": {"lastUsedAt": datetime.now(timezone.utc).isoformat()}
                    },
                    upsert=True
                )
        
    except Exception as e:
        logger.error(f"Claude API error: {str(e)}")
        response_text = f"Error: {str(e)[:100]}"
        citations = []
    
    # Deduplicate citations by source
    unique_citations = {}
    used_sources = []
    for c in citations:
        key = c["sourceId"]
        if key not in unique_citations:
            unique_citations[key] = {
                "sourceName": c["sourceName"],
                "sourceId": c["sourceId"],
                "sourceType": c.get("sourceType", "unknown"),  # project or global
                "chunks": []
            }
            used_sources.append({
                "sourceId": c["sourceId"],
                "sourceName": c["sourceName"],
                "sourceType": c.get("sourceType", "unknown")
            })
        unique_citations[key]["chunks"].append({
            "index": c["chunkIndex"] + 1,
            "chunkId": c.get("chunkId", ""),
            "textFragment": c.get("textFragment", "")
        })
    
    final_citations = list(unique_citations.values()) if unique_citations else None
    final_used_sources = used_sources if used_sources else None
    
    # Save to semantic cache if we have embedding and got a valid response (not from cache and not error)
    if question_embedding and not from_cache and not response_text.startswith("I apologize"):
        await save_to_cache(
            question=message_data.content,
            answer=response_text,
            project_id=project_id,
            embedding=question_embedding,
            user_id=current_user["id"],
            cache_context_hash=cache_context_hash,
            source_ids=active_source_ids,
            sources_used=final_used_sources
        )
    
    # Save assistant message
    assistant_msg_id = str(uuid.uuid4())
    assistant_message = {
        "id": assistant_msg_id,
        "chatId": chat_id,
        "role": "assistant",
        "content": response_text,
        "citations": final_citations,
        "usedSources": final_used_sources,
        "autoIngestedUrls": [s["id"] for s in auto_ingested_sources] if auto_ingested_sources else None,
        "senderEmail": None,
        "senderName": "GPT",
        "fromCache": from_cache,
        "cacheInfo": cache_info,
        "createdAt": datetime.now(timezone.utc).isoformat()
    }
    await db.messages.insert_one(assistant_message)
    
    # Track source usage statistics
    if final_used_sources:
        for source_info in final_used_sources:
            await db.source_usage.update_one(
                {"sourceId": source_info["sourceId"]},
                {
                    "$inc": {"usageCount": 1},
                    "$set": {
                        "lastUsedAt": datetime.now(timezone.utc).isoformat(),
                        "sourceName": source_info["sourceName"]
                    },
                    "$push": {
                        "usageHistory": {
                            "$each": [{
                                "userId": current_user["id"],
                                "userEmail": current_user["email"],
                                "chatId": chat_id,
                                "messageId": assistant_msg_id,
                                "timestamp": datetime.now(timezone.utc).isoformat()
                            }],
                            "$slice": -100  # Keep last 100 uses
                        }
                    }
                },
                upsert=True
            )
    
    return MessageResponse(**assistant_message)

# ==================== SAVE TO KNOWLEDGE ====================

class SaveToKnowledgeRequest(BaseModel):
    content: str
    chatId: Optional[str] = None

@api_router.post("/save-to-knowledge")
async def save_to_knowledge(
    request: SaveToKnowledgeRequest,
    current_user: dict = Depends(get_current_user)
):
    """Save AI message content as a Personal Source"""
    try:
        # Generate source name from content (first 50 chars + timestamp)
        content_preview = request.content[:50].replace('\n', ' ').strip()
        if len(request.content) > 50:
            content_preview += "..."
        timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        source_name = f"{content_preview} ({timestamp})"
        
        # Create source ID
        source_id = str(uuid.uuid4())
        
        # Create the source document (matching enterprise_sources format)
        source_doc = {
            "id": source_id,
            "level": "personal",
            "ownerId": current_user["id"],
            "ownerEmail": current_user["email"],
            "projectId": None,
            "departmentId": None,
            "kind": "knowledge",
            "originalName": source_name,
            "mimeType": "text/plain",
            "sizeBytes": len(request.content.encode('utf-8')),
            "storagePath": None,
            "extractedText": request.content,
            "contentHash": hashlib.sha256(request.content.encode('utf-8')).hexdigest(),
            "status": "active",
            "currentVersion": 1,
            "createdAt": datetime.now(timezone.utc).isoformat(),
            "updatedAt": datetime.now(timezone.utc).isoformat()
        }
        
        await db.sources.insert_one(source_doc)
        
        # Create chunks and embeddings for the saved content
        chunks = chunk_text(request.content, chunk_size=1000)
        
        for i, chunk_text_content in enumerate(chunks):
            try:
                # Generate embedding using OpenAI (keeping existing embedding model)
                embedding_response = openai_client.embeddings.create(
                    model="text-embedding-3-small",
                    input=chunk_text_content
                )
                embedding = embedding_response.data[0].embedding
                
                chunk_doc = {
                    "id": str(uuid.uuid4()),
                    "sourceId": source_id,
                    "sourceName": source_name,
                    "chunkIndex": i,
                    "text": chunk_text_content,
                    "embedding": embedding,
                    "createdAt": datetime.now(timezone.utc).isoformat()
                }
                await db.source_chunks.insert_one(chunk_doc)
            except Exception as e:
                logger.error(f"Error creating embedding for chunk {i}: {str(e)}")
        
        # Log the action
        await db.audit_logs.insert_one({
            "id": str(uuid.uuid4()),
            "userId": current_user["id"],
            "userEmail": current_user["email"],
            "action": "save_to_knowledge",
            "resourceType": "source",
            "resourceId": source_id,
            "details": {"sourceName": source_name, "contentLength": len(request.content)},
            "timestamp": datetime.now(timezone.utc).isoformat()
        })
        
        return {
            "success": True,
            "sourceId": source_id,
            "sourceName": source_name,
            "message": "Saved to Knowledge ✅"
        }
        
    except Exception as e:
        logger.error(f"Error saving to knowledge: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to save: {str(e)}")

# ==================== ADMIN ENDPOINTS ====================

@api_router.get("/admin/source-stats")
async def get_source_stats(current_user: dict = Depends(get_current_user)):
    """Get source statistics per user - for admin dashboard"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Get all users
    users = await db.users.find({}, {"_id": 0, "id": 1, "email": 1}).to_list(1000)
    user_map = {u["id"]: u["email"] for u in users}
    
    # Get all projects with their owners
    projects = await db.projects.find({}, {"_id": 0, "id": 1, "ownerId": 1, "name": 1}).to_list(1000)
    project_owner_map = {p["id"]: p["ownerId"] for p in projects}
    
    # Get all sources with size info
    sources = await db.sources.find({}, {"_id": 0, "projectId": 1, "sizeBytes": 1, "originalName": 1, "createdAt": 1, "kind": 1}).to_list(10000)
    
    # Aggregate by user
    user_stats = {}
    for source in sources:
        project_id = source.get("projectId")
        owner_id = project_owner_map.get(project_id)
        
        if not owner_id:
            continue
            
        if owner_id not in user_stats:
            user_stats[owner_id] = {
                "userId": owner_id,
                "email": user_map.get(owner_id, "Unknown"),
                "sourceCount": 0,
                "totalSizeBytes": 0,
                "fileCount": 0,
                "urlCount": 0
            }
        
        user_stats[owner_id]["sourceCount"] += 1
        user_stats[owner_id]["totalSizeBytes"] += source.get("sizeBytes", 0) or 0
        
        if source.get("kind") == "url":
            user_stats[owner_id]["urlCount"] += 1
        else:
            user_stats[owner_id]["fileCount"] += 1
    
    # Convert to list and sort by size
    result = list(user_stats.values())
    result.sort(key=lambda x: x["totalSizeBytes"], reverse=True)
    
    # Calculate totals
    total_sources = sum(u["sourceCount"] for u in result)
    total_size = sum(u["totalSizeBytes"] for u in result)
    
    return {
        "users": result,
        "totalSources": total_sources,
        "totalSizeBytes": total_size
    }

@api_router.get("/admin/global-sources/stats")
async def get_global_sources_usage_stats(current_user: dict = Depends(get_current_user)):
    """Get usage statistics for all global sources"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Get all global sources
    sources = await db.sources.find(
        {"projectId": GLOBAL_PROJECT_ID}, 
        {"_id": 0, "id": 1, "originalName": 1, "url": 1, "chunkCount": 1, "sizeBytes": 1, "createdAt": 1}
    ).to_list(1000)
    
    # Get usage stats for each source
    result = []
    for source in sources:
        usage = await db.source_usage.find_one(
            {"sourceId": source["id"]}, 
            {"_id": 0}
        )
        
        source_name = source.get("originalName") or source.get("url") or "Unknown"
        
        result.append({
            "sourceId": source["id"],
            "sourceName": source_name,
            "chunkCount": source.get("chunkCount", 0),
            "sizeBytes": source.get("sizeBytes", 0),
            "createdAt": source.get("createdAt"),
            "usageCount": usage.get("usageCount", 0) if usage else 0,
            "lastUsedAt": usage.get("lastUsedAt") if usage else None,
            "recentUsers": [
                {"email": u["userEmail"], "timestamp": u["timestamp"]} 
                for u in (usage.get("usageHistory", []) if usage else [])[-5:]
            ]
        })
    
    # Sort by usage count descending
    result.sort(key=lambda x: x["usageCount"], reverse=True)
    
    # Calculate totals
    total_usage = sum(s["usageCount"] for s in result)
    sources_used = sum(1 for s in result if s["usageCount"] > 0)
    
    return {
        "sources": result,
        "totalUsageCount": total_usage,
        "sourcesUsedCount": sources_used,
        "totalSourcesCount": len(result)
    }

@api_router.get("/admin/cache/stats")
async def get_cache_stats(current_user: dict = Depends(get_current_user)):
    """Get semantic cache statistics"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Get cache entries
    cache_entries = await db.semantic_cache.find({}, {"_id": 0, "embedding": 0}).to_list(1000)
    
    total_entries = len(cache_entries)
    total_hits = sum(e.get("hitCount", 0) for e in cache_entries)
    
    # Group by project
    by_project = {}
    for entry in cache_entries:
        pid = entry.get("projectId") or "global"
        if pid not in by_project:
            by_project[pid] = {"count": 0, "hits": 0}
        by_project[pid]["count"] += 1
        by_project[pid]["hits"] += entry.get("hitCount", 0)
    
    # Get top cached questions
    top_entries = sorted(cache_entries, key=lambda x: x.get("hitCount", 0), reverse=True)[:10]
    
    return {
        "totalEntries": total_entries,
        "totalHits": total_hits,
        "byProject": by_project,
        "topEntries": [{
            "id": e["id"],
            "question": e["question"][:100] + "..." if len(e.get("question", "")) > 100 else e.get("question", ""),
            "hitCount": e.get("hitCount", 0),
            "lastHitAt": e.get("lastHitAt"),
            "createdAt": e.get("createdAt")
        } for e in top_entries],
        "settings": {
            "similarityThreshold": CACHE_SIMILARITY_THRESHOLD,
            "ttlDays": CACHE_TTL_DAYS
        }
    }

@api_router.delete("/admin/cache/clear")
async def clear_cache(current_user: dict = Depends(get_current_user)):
    """Clear all semantic cache entries"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    result = await db.semantic_cache.delete_many({})
    return {"message": f"Cleared {result.deleted_count} cache entries"}

@api_router.delete("/admin/cache/{cache_id}")
async def delete_cache_entry(cache_id: str, current_user: dict = Depends(get_current_user)):
    """Delete specific cache entry"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    result = await db.semantic_cache.delete_one({"id": cache_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Cache entry not found")
    return {"message": "Cache entry deleted"}

# ==================== GLOBAL SOURCES ENDPOINTS ====================

async def can_edit_global_sources(user: dict) -> bool:
    """Check if user can edit global sources"""
    if is_admin(user["email"]):
        return True
    return user.get("canEditGlobalSources", False)

@api_router.get("/admin/global-sources")
async def get_global_sources(current_user: dict = Depends(get_current_user)):
    """Get all global sources - admin only for full list"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    sources = await db.sources.find({"projectId": GLOBAL_PROJECT_ID}, {"_id": 0}).to_list(1000)
    return sources

@api_router.get("/global-sources")
async def get_global_sources_for_users(current_user: dict = Depends(get_current_user)):
    """Get global sources for any authenticated user"""
    sources = await db.sources.find({"projectId": GLOBAL_PROJECT_ID}, {"_id": 0}).to_list(1000)
    # Add canEdit flag for frontend
    can_edit = await can_edit_global_sources(current_user)
    return {"sources": sources, "canEdit": can_edit}

@api_router.post("/global-sources/upload")
async def user_upload_global_source(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    """User with permission uploads a global source file"""
    if not await can_edit_global_sources(current_user):
        raise HTTPException(status_code=403, detail="You don't have permission to edit global sources")
    
    # Reuse the same logic as admin upload
    allowed_types = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "text/plain": "txt",
        "text/markdown": "md",
        "text/csv": "csv",
        "application/csv": "csv",
        "image/png": "png",
        "image/jpeg": "jpeg"
    }
    
    if file.content_type not in allowed_types:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.content_type}")
    
    file_type = allowed_types[file.content_type]
    content = await file.read()
    
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail=f"File too large. Max size: {MAX_FILE_SIZE // (1024*1024)}MB")
    
    source_id = str(uuid.uuid4())
    file_ext = file.filename.split(".")[-1] if "." in file.filename else file_type
    storage_name = f"global_{source_id}.{file_ext}"
    file_path = UPLOAD_DIR / storage_name
    
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
    
    # Extract text based on file type
    if file_type == "pdf":
        extracted_text = extract_text_from_pdf(content)
    elif file_type == "docx":
        extracted_text = extract_text_from_docx(content)
    elif file_type == "pptx":
        extracted_text = extract_text_from_pptx(content)
    elif file_type == "xlsx":
        extracted_text = extract_text_from_xlsx(content)
    elif file_type == "csv":
        extracted_text = extract_text_from_csv(content)
    elif file_type in ["png", "jpeg", "jpg"]:
        extracted_text = extract_text_from_image(content)
    else:
        extracted_text = extract_text_from_txt(content)
    
    chunks = chunk_text(extracted_text)
    
    source_doc = {
        "id": source_id,
        "projectId": GLOBAL_PROJECT_ID,
        "kind": "file",
        "originalName": file.filename,
        "mimeType": file.content_type,
        "storagePath": storage_name,
        "sizeBytes": len(content),
        "uploadedBy": current_user["id"],
        "createdAt": datetime.now(timezone.utc).isoformat(),
        "chunkCount": len(chunks)
    }
    await db.sources.insert_one(source_doc)
    
    for i, chunk_content in enumerate(chunks):
        chunk_doc = {
            "id": str(uuid.uuid4()),
            "sourceId": source_id,
            "projectId": GLOBAL_PROJECT_ID,
            "chunkIndex": i,
            "content": chunk_content,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.source_chunks.insert_one(chunk_doc)
    
    logger.info(f"User {current_user['email']} uploaded global source {file.filename}")
    return {**source_doc, "_id": None}

@api_router.delete("/global-sources/{source_id}")
async def user_delete_global_source(source_id: str, current_user: dict = Depends(get_current_user)):
    """User with permission deletes a global source"""
    if not await can_edit_global_sources(current_user):
        raise HTTPException(status_code=403, detail="You don't have permission to edit global sources")
    
    source = await db.sources.find_one({"id": source_id, "projectId": GLOBAL_PROJECT_ID}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Global source not found")
    
    # Non-admins can only delete their own uploads
    if not is_admin(current_user["email"]) and source.get("uploadedBy") != current_user["id"]:
        raise HTTPException(status_code=403, detail="You can only delete sources you uploaded")
    
    if source.get("storagePath"):
        file_path = UPLOAD_DIR / source["storagePath"]
        if file_path.exists():
            file_path.unlink()
    
    await db.source_chunks.delete_many({"sourceId": source_id})
    await db.sources.delete_one({"id": source_id})
    
    return {"message": "Global source deleted"}

@api_router.get("/global-sources/{source_id}/preview")
async def user_preview_global_source(source_id: str, current_user: dict = Depends(get_current_user)):
    """Preview global source content - any authenticated user"""
    source = await db.sources.find_one({"id": source_id, "projectId": GLOBAL_PROJECT_ID}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Global source not found")
    
    chunks = await db.source_chunks.find(
        {"sourceId": source_id},
        {"_id": 0, "content": 1, "text": 1, "chunkIndex": 1}
    ).sort("chunkIndex", 1).to_list(1000)
    
    full_text = "\n\n".join([c.get("content") or c.get("text", "") for c in chunks])
    
    return {
        "id": source_id,
        "name": source.get("originalName") or source.get("url"),
        "text": full_text,
        "chunkCount": len(chunks),
        "wordCount": len(full_text.split()),
        "uploadedBy": source.get("uploadedBy")
    }

@api_router.put("/admin/users/{user_id}/global-permission")
async def update_user_global_permission(
    user_id: str, 
    data: UpdateUserGlobalPermissionRequest,
    current_user: dict = Depends(get_current_user)
):
    """Admin grants/revokes global source editing permission"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    user = await db.users.find_one({"id": user_id}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    await db.users.update_one(
        {"id": user_id},
        {"$set": {"canEditGlobalSources": data.canEditGlobalSources}}
    )
    
    return {"message": f"Global sources permission {'granted' if data.canEditGlobalSources else 'revoked'}"}

@api_router.post("/admin/global-sources/upload")
async def upload_global_source(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    """Admin uploads a global source file"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Validate file type
    allowed_types = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "text/plain": "txt",
        "text/markdown": "md",
        "text/csv": "csv",
        "application/csv": "csv",
        "image/png": "png",
        "image/jpeg": "jpeg"
    }
    
    if file.content_type not in allowed_types:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.content_type}")
    
    file_type = allowed_types[file.content_type]
    
    # Check file size
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail=f"File too large. Max size: {MAX_FILE_SIZE // (1024*1024)}MB")
    
    # Save file
    source_id = str(uuid.uuid4())
    file_ext = file.filename.split(".")[-1] if "." in file.filename else file_type
    storage_name = f"global_{source_id}.{file_ext}"
    file_path = UPLOAD_DIR / storage_name
    
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
    
    # Extract text based on file type
    if file_type == "pdf":
        extracted_text = extract_text_from_pdf(content)
    elif file_type == "docx":
        extracted_text = extract_text_from_docx(content)
    elif file_type == "pptx":
        extracted_text = extract_text_from_pptx(content)
    elif file_type == "xlsx":
        extracted_text = extract_text_from_xlsx(content)
    elif file_type == "csv":
        extracted_text = extract_text_from_csv(content)
    elif file_type in ["png", "jpeg", "jpg"]:
        extracted_text = extract_text_from_image(content)
    else:
        extracted_text = extract_text_from_txt(content)
    
    # Chunk the text
    chunks = chunk_text(extracted_text)
    
    # Save source
    source_doc = {
        "id": source_id,
        "projectId": GLOBAL_PROJECT_ID,
        "kind": "file",
        "originalName": file.filename,
        "mimeType": file.content_type,
        "storagePath": storage_name,
        "sizeBytes": len(content),
        "createdAt": datetime.now(timezone.utc).isoformat(),
        "chunkCount": len(chunks)
    }
    await db.sources.insert_one(source_doc)
    
    # Save chunks
    for i, chunk_content in enumerate(chunks):
        chunk_doc = {
            "id": str(uuid.uuid4()),
            "sourceId": source_id,
            "projectId": GLOBAL_PROJECT_ID,
            "chunkIndex": i,
            "content": chunk_content,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.source_chunks.insert_one(chunk_doc)
    
    logger.info(f"Uploaded global source {file.filename} with {len(chunks)} chunks")
    
    return {**source_doc, "_id": None}

@api_router.post("/admin/global-sources/url")
async def add_global_url_source(url_data: UrlSourceCreate, current_user: dict = Depends(get_current_user)):
    """Admin adds a URL as global source"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    url = str(url_data.url)
    
    # Fetch URL content
    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
            response = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
            response.raise_for_status()
            html_content = response.text
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {str(e)}")
    
    # Extract text
    extracted_text = extract_text_from_html(html_content)
    
    if not extracted_text or len(extracted_text) < 50:
        raise HTTPException(status_code=400, detail="Could not extract meaningful content from URL")
    
    # Chunk the text
    chunks = chunk_text(extracted_text)
    
    # Save source
    source_id = str(uuid.uuid4())
    source_doc = {
        "id": source_id,
        "projectId": GLOBAL_PROJECT_ID,
        "kind": "url",
        "url": url,
        "createdAt": datetime.now(timezone.utc).isoformat(),
        "chunkCount": len(chunks)
    }
    await db.sources.insert_one(source_doc)
    
    # Save chunks
    for i, chunk_content in enumerate(chunks):
        chunk_doc = {
            "id": str(uuid.uuid4()),
            "sourceId": source_id,
            "projectId": GLOBAL_PROJECT_ID,
            "chunkIndex": i,
            "content": chunk_content,
            "createdAt": datetime.now(timezone.utc).isoformat()
        }
        await db.source_chunks.insert_one(chunk_doc)
    
    return {**source_doc, "_id": None}

@api_router.delete("/admin/global-sources/{source_id}")
async def delete_global_source(source_id: str, current_user: dict = Depends(get_current_user)):
    """Admin deletes a global source"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    source = await db.sources.find_one({"id": source_id, "projectId": GLOBAL_PROJECT_ID}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Global source not found")
    
    # Delete file if exists
    if source.get("storagePath"):
        file_path = UPLOAD_DIR / source["storagePath"]
        if file_path.exists():
            file_path.unlink()
    
    # Delete chunks
    await db.source_chunks.delete_many({"sourceId": source_id})
    
    # Delete source
    await db.sources.delete_one({"id": source_id})
    
    return {"message": "Global source deleted"}

@api_router.get("/admin/global-sources/{source_id}/preview")
async def preview_global_source(source_id: str, current_user: dict = Depends(get_current_user)):
    """Preview global source content"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    source = await db.sources.find_one({"id": source_id, "projectId": GLOBAL_PROJECT_ID}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Global source not found")
    
    chunks = await db.source_chunks.find(
        {"sourceId": source_id},
        {"_id": 0, "content": 1, "text": 1, "chunkIndex": 1}
    ).sort("chunkIndex", 1).to_list(1000)
    
    full_text = "\n\n".join([c.get("content") or c.get("text", "") for c in chunks])
    
    return {
        "id": source_id,
        "name": source.get("originalName") or source.get("url"),
        "text": full_text,
        "chunkCount": len(chunks),
        "wordCount": len(full_text.split())
    }

@api_router.get("/admin/config", response_model=GPTConfigResponse)
async def get_gpt_config(current_user: dict = Depends(get_current_user)):
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    config = await ensure_gpt_config()
    return GPTConfigResponse(**config)

@api_router.put("/admin/config", response_model=GPTConfigResponse)
async def update_gpt_config(config_data: GPTConfigUpdate, current_user: dict = Depends(get_current_user)):
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    await ensure_gpt_config()
    
    update_data = {"updatedAt": datetime.now(timezone.utc).isoformat()}
    if config_data.model is not None:
        update_data["model"] = config_data.model
    if config_data.developerPrompt is not None:
        update_data["developerPrompt"] = config_data.developerPrompt
    
    await db.gpt_config.update_one({"id": "1"}, {"$set": update_data})
    
    updated_config = await db.gpt_config.find_one({"id": "1"}, {"_id": 0})
    return GPTConfigResponse(**updated_config)

# ==================== ADMIN USER MANAGEMENT ====================

@api_router.post("/admin/users", response_model=UserResponse)
async def admin_create_user(user_data: UserCreate, current_user: dict = Depends(get_current_user)):
    """Admin creates a new user"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    existing = await db.users.find_one({"email": user_data.email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    user_id = str(uuid.uuid4())
    user = {
        "id": user_id,
        "email": user_data.email,
        "passwordHash": hash_password(user_data.password),
        "createdAt": datetime.now(timezone.utc).isoformat(),
        # Enterprise Knowledge Architecture fields
        "departments": [],  # List of department IDs user belongs to
        "primaryDepartmentId": None,  # Default department for the user
        "canEditGlobalSources": False
    }
    await db.users.insert_one(user)
    
    return UserResponse(
        id=user_id,
        email=user_data.email,
        isAdmin=is_admin(user_data.email),
        createdAt=user["createdAt"]
    )

@api_router.get("/users/list")
async def list_users_for_sharing(current_user: dict = Depends(get_current_user)):
    """Get list of all users (for sharing projects)"""
    users = await db.users.find({}, {"_id": 0, "passwordHash": 0}).to_list(1000)
    # Exclude current user and return only id and email
    return [{"id": u["id"], "email": u["email"]} for u in users if u["id"] != current_user["id"]]

@api_router.put("/users/me/primary-department")
async def set_primary_department(
    data: dict,
    current_user: dict = Depends(get_current_user)
):
    """Set user's primary department"""
    department_id = data.get("departmentId")
    
    if department_id:
        # Verify department exists and user is member
        department = await db.departments.find_one({"id": department_id}, {"_id": 0})
        if not department:
            raise HTTPException(status_code=404, detail="Department not found")
        
        user_depts = current_user.get("departments", [])
        if department_id not in user_depts:
            raise HTTPException(status_code=403, detail="You are not a member of this department")
    
    await db.users.update_one(
        {"id": current_user["id"]},
        {"$set": {"primaryDepartmentId": department_id}}
    )
    
    return {"message": "Primary department updated", "primaryDepartmentId": department_id}

@api_router.get("/users/me/departments")
async def get_my_departments(current_user: dict = Depends(get_current_user)):
    """Get current user's departments with details"""
    user_dept_ids = current_user.get("departments", [])
    
    if not user_dept_ids:
        return []
    
    departments = await db.departments.find(
        {"id": {"$in": user_dept_ids}},
        {"_id": 0}
    ).to_list(100)
    
    # Add isManager flag
    for dept in departments:
        dept["isManager"] = current_user["id"] in dept.get("managers", [])
        dept["isPrimary"] = dept["id"] == current_user.get("primaryDepartmentId")
    
    return departments

@api_router.get("/admin/users", response_model=List[UserWithUsageResponse])
async def admin_list_users(current_user: dict = Depends(get_current_user)):
    """Admin gets list of all users with token usage"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    users = await db.users.find({}, {"_id": 0, "passwordHash": 0}).to_list(1000)
    
    result = []
    for user in users:
        # Get token usage for this user
        usage = await db.token_usage.find_one({"userId": user["id"]}, {"_id": 0})
        total_tokens = usage.get("totalTokens", 0) if usage else 0
        message_count = usage.get("messageCount", 0) if usage else 0
        
        result.append(UserWithUsageResponse(
            id=user["id"],
            email=user["email"],
            isAdmin=is_admin(user["email"]),
            createdAt=user["createdAt"],
            totalTokensUsed=total_tokens,
            totalMessagesCount=message_count,
            canEditGlobalSources=user.get("canEditGlobalSources", False)
        ))
    
    return result

@api_router.delete("/admin/users/{user_id}")
async def admin_delete_user(user_id: str, current_user: dict = Depends(get_current_user)):
    """Admin deletes a user"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    user = await db.users.find_one({"id": user_id})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Don't allow deleting yourself
    if user_id == current_user["id"]:
        raise HTTPException(status_code=400, detail="Cannot delete yourself")
    
    # Delete user's data
    await db.users.delete_one({"id": user_id})
    await db.token_usage.delete_one({"userId": user_id})
    await db.user_prompts.delete_one({"userId": user_id})
    
    # Delete user's quick chats
    await db.chats.delete_many({"ownerId": user_id})
    
    # Delete user's projects and their data
    projects = await db.projects.find({"ownerId": user_id}).to_list(1000)
    for project in projects:
        project_id = project["id"]
        chats = await db.chats.find({"projectId": project_id}).to_list(1000)
        for chat in chats:
            await db.messages.delete_many({"chatId": chat["id"]})
        await db.chats.delete_many({"projectId": project_id})
        await db.sources.delete_many({"projectId": project_id})
        await db.source_chunks.delete_many({"projectId": project_id})
        await db.generated_images.delete_many({"projectId": project_id})
    await db.projects.delete_many({"ownerId": user_id})
    
    return {"message": "User deleted successfully"}

@api_router.get("/admin/users/{user_id}/details")
async def get_user_details(user_id: str, current_user: dict = Depends(get_current_user)):
    """Get detailed user info for admin - prompt, projects, activity"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Get user
    user = await db.users.find_one({"id": user_id}, {"_id": 0, "passwordHash": 0})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Get user's prompt
    user_prompt = await db.user_prompts.find_one({"userId": user_id}, {"_id": 0})
    
    # Get user's projects
    projects = await db.projects.find({"ownerId": user_id}, {"_id": 0}).to_list(100)
    
    # Get project stats
    projects_with_stats = []
    for p in projects:
        chat_count = await db.chats.count_documents({"projectId": p["id"]})
        source_count = await db.sources.count_documents({"projectId": p["id"]})
        projects_with_stats.append({
            **p,
            "chatCount": chat_count,
            "sourceCount": source_count
        })
    
    # Get token usage
    usage = await db.token_usage.find_one({"userId": user_id}, {"_id": 0})
    
    # Get recent activity (last 20 messages)
    user_messages = await db.messages.find(
        {"senderEmail": user["email"], "role": "user"},
        {"_id": 0, "id": 1, "chatId": 1, "content": 1, "createdAt": 1}
    ).sort("createdAt", -1).to_list(20)
    
    # Get user-specific model if set
    user_model = user.get("gptModel")
    
    # Get prompt - check both field names for compatibility
    prompt_text = ""
    if user_prompt:
        prompt_text = user_prompt.get("customPrompt") or user_prompt.get("prompt", "")
    
    return {
        "user": user,
        "prompt": prompt_text,
        "gptModel": user_model,
        "projects": projects_with_stats,
        "tokenUsage": {
            "totalTokens": usage.get("totalTokensUsed", 0) if usage else 0,
            "totalMessages": usage.get("totalMessagesCount", 0) if usage else 0
        },
        "recentActivity": user_messages
    }

@api_router.put("/admin/users/{user_id}/prompt")
async def update_user_prompt_admin(user_id: str, data: UserPromptUpdate, current_user: dict = Depends(get_current_user)):
    """Admin updates user's custom prompt"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Check user exists
    user = await db.users.find_one({"id": user_id}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Update or create prompt - use customPrompt field for consistency
    await db.user_prompts.update_one(
        {"userId": user_id},
        {"$set": {"userId": user_id, "customPrompt": data.prompt, "updatedAt": datetime.now(timezone.utc).isoformat()}},
        upsert=True
    )
    
    return {"message": "Prompt updated"}

@api_router.put("/admin/users/{user_id}/model")
async def update_user_model(user_id: str, current_user: dict = Depends(get_current_user), model: str = None):
    """Admin sets user-specific GPT model (overrides global)"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    # Check user exists
    user = await db.users.find_one({"id": user_id}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Update user's model (None = use global)
    await db.users.update_one(
        {"id": user_id},
        {"$set": {"gptModel": model}}
    )
    
    return {"message": "Model updated"}

class UpdateUserModelRequest(BaseModel):
    model: Optional[str] = None

@api_router.put("/admin/users/{user_id}/gpt-model")
async def update_user_gpt_model(user_id: str, data: UpdateUserModelRequest, current_user: dict = Depends(get_current_user)):
    """Admin sets user-specific GPT model (overrides global)"""
    if not is_admin(current_user["email"]):
        raise HTTPException(status_code=403, detail="Admin access required")
    
    user = await db.users.find_one({"id": user_id}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    await db.users.update_one(
        {"id": user_id},
        {"$set": {"gptModel": data.model}}
    )
    
    return {"message": "Model updated", "model": data.model}

# ==================== USER PROMPT ENDPOINTS ====================

@api_router.get("/user/prompt", response_model=UserPromptResponse)
async def get_user_prompt(current_user: dict = Depends(get_current_user)):
    """Get the current user's custom GPT prompt"""
    user_prompt = await db.user_prompts.find_one({"userId": current_user["id"]}, {"_id": 0})
    
    if not user_prompt:
        return UserPromptResponse(
            userId=current_user["id"],
            customPrompt=None,
            updatedAt=datetime.now(timezone.utc).isoformat()
        )
    
    return UserPromptResponse(**user_prompt)

@api_router.put("/user/prompt", response_model=UserPromptResponse)
async def update_user_prompt(data: UserPromptUpdate, current_user: dict = Depends(get_current_user)):
    """Update the current user's custom GPT prompt"""
    now = datetime.now(timezone.utc).isoformat()
    
    existing = await db.user_prompts.find_one({"userId": current_user["id"]})
    
    if existing:
        await db.user_prompts.update_one(
            {"userId": current_user["id"]},
            {"$set": {"customPrompt": data.customPrompt, "updatedAt": now}}
        )
    else:
        await db.user_prompts.insert_one({
            "userId": current_user["id"],
            "customPrompt": data.customPrompt,
            "updatedAt": now
        })
    
    return UserPromptResponse(
        userId=current_user["id"],
        customPrompt=data.customPrompt,
        updatedAt=now
    )

# ==================== HEALTH CHECK ====================

@api_router.get("/")
async def root():
    return {"message": "Shared Project GPT API is running"}

@api_router.get("/health")
async def health():
    return {"status": "healthy", "timestamp": datetime.now(timezone.utc).isoformat()}

# ==================== ENTERPRISE KNOWLEDGE ARCHITECTURE ====================

# Initialize enterprise services
audit_service = AuditService(db)
version_service = VersionService(db)
hierarchical_retrieval = HierarchicalRetrieval(db)

# Helper function for text extraction (wrapper for existing functions)
async def extract_text_wrapper(content: bytes, file_type: str) -> str:
    """Wrapper to extract text based on file type"""
    if file_type == "pdf":
        return extract_text_from_pdf(content)
    elif file_type == "docx":
        return extract_text_from_docx(content)
    elif file_type == "pptx":
        return extract_text_from_pptx(content)
    elif file_type == "xlsx":
        return extract_text_from_xlsx(content)
    elif file_type == "csv":
        return extract_text_from_csv(content)
    elif file_type in ["png", "jpeg", "jpg"]:
        return extract_text_from_image(content)
    else:
        return extract_text_from_txt(content)

# ==================== SOURCE INSIGHTS & SMART QUESTIONS ====================

class SourceInsightsResponse(BaseModel):
    summary: str
    suggestedQuestions: List[str]
    generatedAt: str

class SaveInsightsRequest(BaseModel):
    summary: str
    suggestedQuestions: List[str]

class SmartQuestionsResponse(BaseModel):
    questions: List[str]
    sourceNames: List[str]
    generatedAt: str

@api_router.post("/sources/{source_id}/analyze", response_model=SourceInsightsResponse)
async def analyze_source(source_id: str, current_user: dict = Depends(get_current_user)):
    """
    Analyze a source and generate insights (summary + suggested questions).
    Available to ALL users who can see the source.
    """
    # Find source in any collection
    source = await db.sources.find_one({"id": source_id}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    
    # Get source chunks
    chunks = await db.source_chunks.find(
        {"sourceId": source_id},
        {"_id": 0}
    ).sort("chunkIndex", 1).to_list(50)
    
    if not chunks:
        raise HTTPException(status_code=400, detail="Source has no content to analyze")
    
    # Combine chunks (handle both 'content' and 'text' fields)
    full_text = "\n\n".join([c.get("content") or c.get("text", "") for c in chunks])
    
    # Limit text for analysis
    text_for_analysis = full_text[:8000]
    source_name = source.get("originalName") or source.get("url") or "Unknown"
    
    # Use Claude for analysis
    try:
        import anthropic
        CLAUDE_API_KEY = os.environ.get('CLAUDE_API_KEY', '')
        claude_client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)
        
        analysis_prompt = f"""Analyze the following document and provide:
1. A brief summary (2-3 sentences) describing what this document contains
2. Exactly 5 specific questions that a user could ask about this document's content

IMPORTANT: Detect the language of the document and respond in THE SAME LANGUAGE as the document content. 
- If the document is in Armenian, respond in Armenian
- If the document is in Russian, respond in Russian  
- If the document is in English, respond in English
- etc.

Document name: {source_name}
Document content:
{text_for_analysis}

Respond in JSON format:
{{
  "summary": "Your 2-3 sentence summary here (in document's language)",
  "questions": ["Question 1?", "Question 2?", "Question 3?", "Question 4?", "Question 5?"]
}}

Important: Respond ONLY with valid JSON, no additional text. Questions and summary must be in the SAME LANGUAGE as the document."""

        response = claude_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1024,
            messages=[{"role": "user", "content": analysis_prompt}]
        )
        
        import json
        result_text = response.content[0].text.strip()
        # Clean up potential markdown code blocks
        if result_text.startswith("```"):
            result_text = result_text.split("```")[1]
            if result_text.startswith("json"):
                result_text = result_text[4:]
        result_text = result_text.strip()
        
        result = json.loads(result_text)
        
        return SourceInsightsResponse(
            summary=result.get("summary", "Unable to generate summary"),
            suggestedQuestions=result.get("questions", [])[:5],
            generatedAt=datetime.now(timezone.utc).isoformat()
        )
        
    except Exception as e:
        logger.error(f"Source analysis error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)[:100]}")


@api_router.post("/sources/{source_id}/save-insights")
async def save_source_insights(
    source_id: str, 
    data: SaveInsightsRequest,
    current_user: dict = Depends(get_current_user)
):
    """Save generated insights to a source"""
    source = await db.sources.find_one({"id": source_id}, {"_id": 0})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    
    await db.sources.update_one(
        {"id": source_id},
        {"$set": {
            "insights": {
                "summary": data.summary,
                "suggestedQuestions": data.suggestedQuestions,
                "savedAt": datetime.now(timezone.utc).isoformat(),
                "savedBy": current_user["id"]
            }
        }}
    )
    
    return {"message": "Insights saved successfully"}


@api_router.get("/sources/{source_id}/insights")
async def get_source_insights(source_id: str, current_user: dict = Depends(get_current_user)):
    """Get saved insights for a source"""
    source = await db.sources.find_one({"id": source_id}, {"_id": 0, "insights": 1})
    if not source:
        raise HTTPException(status_code=404, detail="Source not found")
    
    insights = source.get("insights")
    if not insights:
        return {"hasInsights": False}
    
    return {
        "hasInsights": True,
        "summary": insights.get("summary"),
        "suggestedQuestions": insights.get("suggestedQuestions", []),
        "savedAt": insights.get("savedAt")
    }


@api_router.post("/chats/{chat_id}/smart-questions", response_model=SmartQuestionsResponse)
async def generate_smart_questions(chat_id: str, current_user: dict = Depends(get_current_user)):
    """
    Generate smart question suggestions based on active sources in the chat.
    """
    chat = await db.chats.find_one({"id": chat_id}, {"_id": 0})
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    
    project_id = chat.get("projectId")
    source_mode = chat.get("sourceMode", "all")
    
    # Gather active source IDs based on source mode
    active_source_ids = []
    user_department_ids = current_user.get("departments", [])
    
    # Personal sources
    personal_sources = await db.sources.find({
        "level": "personal",
        "ownerId": current_user["id"],
        "status": {"$in": ["active", None]}
    }, {"_id": 0, "id": 1}).to_list(100)
    active_source_ids.extend([s["id"] for s in personal_sources])
    
    # Project sources
    if project_id:
        project_sources = await db.sources.find({
            "projectId": project_id,
            "level": {"$in": ["project", None]},
            "status": {"$in": ["active", None]}
        }, {"_id": 0, "id": 1}).to_list(100)
        active_source_ids.extend([s["id"] for s in project_sources])
    
    # Department and global sources only if mode is 'all'
    if source_mode == 'all':
        if user_department_ids:
            dept_sources = await db.sources.find({
                "departmentId": {"$in": user_department_ids},
                "level": "department",
                "status": "active"
            }, {"_id": 0, "id": 1}).to_list(100)
            active_source_ids.extend([s["id"] for s in dept_sources])
        
        global_sources = await db.sources.find({
            "$or": [
                {"projectId": "__global__"},
                {"level": "global", "status": "active"}
            ]
        }, {"_id": 0, "id": 1}).to_list(100)
        active_source_ids.extend([s["id"] for s in global_sources])
    
    if not active_source_ids:
        raise HTTPException(status_code=400, detail="No active sources available")
    
    # Get source names and sample content
    sources = await db.sources.find(
        {"id": {"$in": active_source_ids}},
        {"_id": 0, "id": 1, "originalName": 1, "url": 1}
    ).to_list(100)
    
    source_names = [s.get("originalName") or s.get("url") or "Unknown" for s in sources]
    
    # Get sample chunks from each source (limit to prevent token overflow)
    sample_content = []
    for source in sources[:5]:  # Limit to 5 sources
        chunks = await db.source_chunks.find(
            {"sourceId": source["id"]},
            {"_id": 0}
        ).sort("chunkIndex", 1).to_list(3)  # 3 chunks per source
        
        source_name = source.get("originalName") or source.get("url") or "Unknown"
        for chunk in chunks:
            text = chunk.get("content") or chunk.get("text", "")
            sample_content.append(f"[{source_name}]: {text[:500]}")
    
    combined_content = "\n\n".join(sample_content)[:6000]
    
    # Generate questions using Claude
    try:
        import anthropic
        CLAUDE_API_KEY = os.environ.get('CLAUDE_API_KEY', '')
        claude_client = anthropic.Anthropic(api_key=CLAUDE_API_KEY)
        
        prompt = f"""Based on the following document excerpts, generate exactly 5 specific, useful questions that a user might want to ask about this content.

IMPORTANT: Detect the language of the content and generate questions in THE SAME LANGUAGE.
- If content is in Armenian, questions must be in Armenian
- If content is in Russian, questions must be in Russian
- If content is in English, questions must be in English

Available sources: {', '.join(source_names[:10])}

Content excerpts:
{combined_content}

Generate 5 practical questions that:
- Are specific to the actual content shown
- Would be useful for someone working with these documents
- Cover different aspects of the content
- Are in the SAME LANGUAGE as the content

Respond with ONLY a JSON array of 5 questions:
["Question 1?", "Question 2?", "Question 3?", "Question 4?", "Question 5?"]"""

        response = claude_client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=512,
            messages=[{"role": "user", "content": prompt}]
        )
        
        import json
        result_text = response.content[0].text.strip()
        # Clean up potential markdown
        if result_text.startswith("```"):
            result_text = result_text.split("```")[1]
            if result_text.startswith("json"):
                result_text = result_text[4:]
        result_text = result_text.strip()
        
        questions = json.loads(result_text)
        
        return SmartQuestionsResponse(
            questions=questions[:5],
            sourceNames=source_names[:5],
            generatedAt=datetime.now(timezone.utc).isoformat()
        )
        
    except Exception as e:
        logger.error(f"Smart questions error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to generate questions: {str(e)[:100]}")


# Setup enterprise routes with dependencies
setup_department_routes(db, get_current_user, is_admin, audit_service)
setup_enterprise_source_routes(
    db, 
    get_current_user, 
    is_admin, 
    audit_service, 
    version_service,
    extract_text_wrapper,
    chunk_text,
    chunk_tabular_text,
    UPLOAD_DIR,
    MAX_FILE_SIZE,
    SUPPORTED_MIME_TYPES
)

# Include the router in the main app
app.include_router(api_router)
app.include_router(departments_router)
app.include_router(enterprise_sources_router)
app.include_router(news_router, prefix="/api")

# Setup analyzer routes
setup_analyzer_routes(db, get_current_user)
app.include_router(analyzer_router, prefix="/api")

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
