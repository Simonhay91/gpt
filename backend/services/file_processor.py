"""Text extraction utilities for various file types"""
import io
import logging
from typing import List
from fastapi import HTTPException
from docx import Document
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup

_markitdown = None

def _get_markitdown():
    global _markitdown
    if _markitdown is None:
        try:
            from markitdown import MarkItDown
            _markitdown = MarkItDown()
        except Exception:
            pass
    return _markitdown

logger = logging.getLogger(__name__)

# Chunk settings
CHUNK_SIZE = 1500
CHUNK_SIZE_TABULAR = 800
CHUNK_OVERLAP = 200


def extract_text_from_pdf(file_content: bytes) -> str:
    try:
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name
        try:
            md = _get_markitdown()
            if md:
                result = md.convert(tmp_path)
                text = result.text_content.strip()
                if text:
                    return text
        except Exception:
            pass
        finally:
            os.unlink(tmp_path)
        import pdfplumber
        text_parts = []
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts).strip()
    except Exception as e:
        logger.error(f"PDF extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    try:
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_content)
            tmp_path = tmp.name
        try:
            md = _get_markitdown()
            if md:
                result = md.convert(tmp_path)
                text = result.text_content.strip()
                if text:
                    return text
        except Exception:
            pass
        finally:
            os.unlink(tmp_path)
        doc = Document(io.BytesIO(file_content))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.error(f"DOCX extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_txt(file_content: bytes) -> str:
    """Extract text from TXT/MD file content"""
    try:
        return file_content.decode('utf-8')
    except UnicodeDecodeError:
        try:
            return file_content.decode('latin-1')
        except Exception as e:
            logger.error(f"TXT extraction error: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Failed to read text file: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint file content"""
    try:
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation(BytesIO(file_content))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from PowerPoint: {str(e)}")


def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from Excel file content - optimized for AI search"""
    try:
        from openpyxl import load_workbook
        from io import BytesIO

        wb = load_workbook(BytesIO(file_content), data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]

            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            headers = [str(cell).strip() if cell is not None else f"Column_{i}" for i, cell in enumerate(rows[0])]

            for row_idx, row in enumerate(rows[1:], start=2):
                row_values = [str(cell).strip() if cell is not None else "" for cell in row]

                if not any(row_values):
                    continue

                record_parts = []
                for header, value in zip(headers, row_values):
                    if value:
                        record_parts.append(f"{header}: {value}")

                if record_parts:
                    record = f"Row {row_idx}: " + ", ".join(record_parts)
                    sheet_text.append(record)

            if len(sheet_text) > 1:
                sheet_text.insert(1, f"Columns: {', '.join(headers)}")
                text_parts.append("\n".join(sheet_text))

        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"XLSX extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from Excel: {str(e)}")


def extract_text_from_csv(file_content: bytes) -> str:
    """Extract text from CSV file content"""
    import csv
    from io import StringIO

    csv.field_size_limit(10 * 1024 * 1024)

    try:
        text_content = None
        for encoding in ['utf-8', 'utf-8-sig', 'cp1251', 'latin-1']:
            try:
                text_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text_content is None:
            text_content = file_content.decode('utf-8', errors='replace')

        text_content = text_content.replace('\r\n', '\n').replace('\r', '\n')

        reader = csv.reader(StringIO(text_content, newline=''))
        rows = list(reader)

        if not rows:
            return ""

        headers = rows[0] if rows else []

        text_parts = []
        text_parts.append(f"[CSV Data - {len(rows)-1} records]")
        text_parts.append(f"Columns: {', '.join(headers)}")
        text_parts.append("")

        for i, row in enumerate(rows[1:], 1):
            if not any(cell.strip() for cell in row if cell):
                continue

            record_parts = []
            for j, cell in enumerate(row):
                if cell and cell.strip():
                    header = headers[j] if j < len(headers) else f"Column{j+1}"
                    clean_cell = cell.strip().replace('\n', ' ').replace('\r', ' ')
                    record_parts.append(f"{header}: {clean_cell}")

            if record_parts:
                text_parts.append(f"[Record {i}]")
                text_parts.append(" | ".join(record_parts))

        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"CSV extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from CSV: {str(e)}")


def extract_text_from_image(file_content: bytes) -> str:
    """Extract text from image using OCR (pytesseract)"""
    try:
        pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

        image = Image.open(io.BytesIO(file_content))

        if image.mode in ('RGBA', 'P'):
            image = image.convert('RGB')

        text = pytesseract.image_to_string(image, lang='rus+eng')
        text = text.strip()

        if not text:
            return "[Image: No text detected]"

        return f"[Image OCR Content]\n{text}"
    except Exception as e:
        logger.error(f"Image OCR error: {str(e)}")
        return f"[Image: OCR failed - {str(e)[:50]}]"


def extract_text_from_html(html_content: str) -> str:
    """Extract readable text from HTML content"""
    try:
        soup = BeautifulSoup(html_content, 'html.parser')

        for element in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
            element.decompose()

        text = soup.get_text(separator='\n')

        lines = [line.strip() for line in text.splitlines()]
        lines = [line for line in lines if line]

        return '\n\n'.join(lines)
    except Exception as e:
        logger.error(f"HTML extraction error: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to extract text from URL: {str(e)}")


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE) -> List[str]:
    """Split text into overlapping chunks"""
    if not text:
        return []

    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]

        # Paragraph boundary-ով կտրել
        if end < len(text):
            last_para = chunk.rfind('\n\n')
            last_newline = chunk.rfind('\n')
            cut = last_para if last_para > chunk_size * 0.5 else last_newline
            if cut > 0:
                chunk = text[start:start + cut]
                end = start + cut

        chunks.append(chunk.strip())
        start = end - CHUNK_OVERLAP

    return [c for c in chunks if c]


def chunk_tabular_text(text: str, chunk_size: int = CHUNK_SIZE_TABULAR) -> List[str]:
    """Split tabular text (Excel, CSV) into chunks - preserves row integrity"""
    if not text:
        return []

    chunks = []
    current_chunk = ""

    lines = text.split('\n')

    for line in lines:
        if line.startswith('[Sheet:') or line.startswith('Columns:'):
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = line
        elif len(current_chunk) + len(line) + 1 <= chunk_size:
            if current_chunk:
                current_chunk += "\n" + line
            else:
                current_chunk = line
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = line

    if current_chunk:
        chunks.append(current_chunk.strip())

    return chunks